import base64
import io
import inspect
import json
import shutil
import sqlite3
import subprocess
import sys
import urllib.error
import urllib.request
import uuid
from pathlib import Path


ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "webui"))

import server  # noqa: E402
import measure_context_limits  # noqa: E402
from rag.index import build_code_graph_context, code_graph_status, index_is_stale, rebuild_index, search_code_graph, search_index  # noqa: E402


def assert_true(condition, message):
    if not condition:
        raise AssertionError(message)


def test_no_context_chat_payload():
    payload = server.build_raw_chat_user_message("", "hello")
    assert_true("PINNED FILE CONTENT" not in payload, "empty chat context must not mention pinned files")
    assert_true(payload == "USER QUESTION:\nhello", "empty chat context should be plain user question")


def test_request_max_tokens_clamps_to_default():
    assert_true(server.get_request_max_tokens({}, 128) == 128, "missing maxTokens should use default")
    assert_true(server.get_request_max_tokens({"maxTokens": 32}, 128) == 32, "maxTokens should allow smaller smoke-test budget")
    assert_true(server.get_request_max_tokens({"maxTokens": 9999}, 128) == 128, "maxTokens should not exceed server default")
    assert_true(server.get_request_max_tokens({"maxTokens": "bad"}, 128) == 128, "invalid maxTokens should use default")


def test_qwen_request_options_disable_thinking():
    assert_true(
        server.get_model_request_options("qwen36a3b") == {"chat_template_kwargs": {"enable_thinking": False}},
        "qwen36a3b requests should disable thinking for visible CodeWorker answers",
    )
    assert_true(
        server.get_model_request_options("gemma4") == {},
        "non-Qwen requests should not get Qwen chat template options",
    )


def test_default_model_uses_last_used_preference():
    old_path = server.MODEL_PREFERENCES_PATH
    old_model_key = server.STATE.model_key
    old_model_alias = server.STATE.model_alias
    temp_dir = ROOT / ".tmp" / f"regression-model-preferences-{uuid.uuid4().hex}"
    shutil.rmtree(temp_dir, ignore_errors=True)
    temp_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.MODEL_PREFERENCES_PATH = temp_dir / "model-preferences.json"
        assert_true(server.DEFAULT_MODEL_KEY == "gemma4", "static fallback model should remain gemma4")
        assert_true(server.get_default_model_key() == "gemma4", "missing preference should fall back to gemma4")
        assert_true(server.get_models_payload()["defaultModelKey"] == "gemma4", "/api/models should expose fallback default")

        server.save_last_used_model_key("qwen36a3b")
        assert_true(server.get_default_model_key() == "qwen36a3b", "last-used preference should become runtime default")
        assert_true(server.get_models_payload()["defaultModelKey"] == "qwen36a3b", "/api/models should expose last-used default")
        assert_true(server.get_models_payload()["lastUsedModelKey"] == "qwen36a3b", "/api/models should expose last-used model")

        server.clear_session()
        assert_true(server.STATE.model_key == "qwen36a3b", "clear_session should reset to last-used model")
        assert_true(server.STATE.model_alias == server.get_model_alias("qwen36a3b"), "clear_session should update last-used alias")

        server.MODEL_PREFERENCES_PATH.write_text(json.dumps({"lastUsedModelKey": "not-a-model"}, ensure_ascii=False), encoding="utf-8")
        assert_true(server.get_default_model_key() == "gemma4", "invalid preference should fall back to static default")
    finally:
        server.MODEL_PREFERENCES_PATH = old_path
        server.STATE.model_key = old_model_key
        server.STATE.model_alias = old_model_alias
        shutil.rmtree(temp_dir, ignore_errors=True)


def test_chat_exchange_persists_last_used_model_preference():
    old_path = server.MODEL_PREFERENCES_PATH
    old_threads_dir = server.THREADS_DIR
    old_active_thread_id = server.ACTIVE_THREAD_ID
    old_model_key = server.STATE.model_key
    old_model_alias = server.STATE.model_alias
    old_history = list(server.STATE.history)
    old_transcript = list(server.STATE.transcript)
    temp_dir = ROOT / ".tmp" / f"regression-chat-model-preference-{uuid.uuid4().hex}"
    shutil.rmtree(temp_dir, ignore_errors=True)
    temp_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.MODEL_PREFERENCES_PATH = temp_dir / "model-preferences.json"
        server.THREADS_DIR = temp_dir / "chat-threads"
        server.ACTIVE_THREAD_ID = None
        with server.STATE_LOCK:
            server.append_chat_exchange_locked("qwen36a3b", "ping", [], "", "pong")
        assert_true(server.get_default_model_key() == "qwen36a3b", "successful chat exchange should persist last-used model")
    finally:
        server.MODEL_PREFERENCES_PATH = old_path
        server.THREADS_DIR = old_threads_dir
        server.ACTIVE_THREAD_ID = old_active_thread_id
        server.STATE.model_key = old_model_key
        server.STATE.model_alias = old_model_alias
        server.STATE.history = old_history
        server.STATE.transcript = old_transcript
        shutil.rmtree(temp_dir, ignore_errors=True)


def test_edit_plan_flow_uses_requested_model_key():
    source = inspect.getsource(server.WebUIHandler.handle_edit_plan)
    assert_true('requested_model_key = str(payload.get("modelKey", "")).strip().lower()' in source, "edit plan should read the selected model from the request")
    assert_true("model_key = requested_model_key or STATE.model_key" in source, "edit plan should prefer the requested model over stale server state")
    assert_true("ensure_local_model_server(model_key" in source, "edit plan should ensure a newly selected model before planning")
    assert_true("remember_last_used_model_key(snapshot.model_key)" in source, "edit plan should persist the model it actually used")


def test_gemma_context_window_matches_local_bench():
    assert_true(server.get_model_context_limit("gemma4") == 262144, "gemma4 should default to the selectable 256k context window")
    assert_true(server.get_model_context_limit("qwen35") == 262144, "qwen35 should default to the selectable 256k context window")
    assert_true(any(item["value"] == 262144 for item in server.get_context_options_payload()), "context options should expose 256k")
    assert_true(server.get_chat_max_tokens("gemma4") <= 4096, "gemma4 response budget should leave room for input context")
    limits = server.get_context_limits("gemma4", single_file_focus=False)
    assert_true(limits["total_chars"] >= 20000, "gemma4 RAG char budget should use the selected context window")


def test_context_calibration_overrides_input_budget_and_model_payload():
    old_path = server.MODEL_CONTEXT_CALIBRATION_PATH
    old_context = server.MODEL_CONTEXT_SELECTIONS_PATH
    temp_dir = ROOT / ".tmp" / f"regression-context-calibration-{uuid.uuid4().hex}"
    shutil.rmtree(temp_dir, ignore_errors=True)
    temp_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.MODEL_CONTEXT_CALIBRATION_PATH = temp_dir / "model-context-calibration.json"
        server.MODEL_CONTEXT_SELECTIONS_PATH = temp_dir / "model-contexts.json"
        server.MODEL_CONTEXT_CALIBRATION_PATH.write_text(
            json.dumps(
                {
                    "qwen36a3b": {
                        "contextWindow": 32768,
                        "maxInputChars": 41000,
                        "structuredEditChars": 37000,
                        "measuredAt": "2026-05-23T20:00:00+08:00",
                    }
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        assert_true(
            server.estimate_input_char_budget("qwen36a3b", 900) == 37000,
            "structured edit calibration should override conservative context budget",
        )
        limits = server.get_context_limits("qwen36a3b", single_file_focus=True, max_response_tokens=900)
        assert_true(limits["total_chars"] == 37000, "edit context should use calibrated structured edit budget")
        assert_true(limits["full_total_chars"] == 37000, "full-file edit context should expand to calibrated capacity")
        model_payload = server.get_models_payload()["models"]["qwen36a3b"]
        assert_true(model_payload["calibrated"] is True, "/api/models should expose calibration state")
        assert_true(model_payload["structuredEditChars"] == 37000, "/api/models should expose calibrated edit budget")
        assert_true(model_payload["measuredAt"] == "2026-05-23T20:00:00+08:00", "/api/models should expose calibration timestamp")
    finally:
        server.MODEL_CONTEXT_CALIBRATION_PATH = old_path
        server.MODEL_CONTEXT_SELECTIONS_PATH = old_context
        shutil.rmtree(temp_dir, ignore_errors=True)


def test_context_calibration_runner_targets_selected_model_and_context():
    old_calibration_path = server.MODEL_CONTEXT_CALIBRATION_PATH
    old_selection_path = server.MODEL_CONTEXT_SELECTIONS_PATH
    old_logs_dir = server.LOGS_DIR
    old_runner = server.run_python_script_via_log
    temp_dir = ROOT / ".tmp" / f"regression-context-calibration-runner-{uuid.uuid4().hex}"
    shutil.rmtree(temp_dir, ignore_errors=True)
    temp_dir.mkdir(parents=True, exist_ok=True)
    captured = {}

    def fake_runner(script_name, *args, timeout_seconds=None):
        captured["script"] = script_name
        captured["args"] = list(args)
        captured["timeout"] = timeout_seconds
        server.LOGS_DIR.mkdir(parents=True, exist_ok=True)
        (server.LOGS_DIR / "model-context-bench.json").write_text(
            json.dumps(
                [
                    {
                        "model": "qwen36a3b",
                        "context": 32768,
                        "startup_ok": True,
                        "tests": {
                            "entry": {"ok": True},
                            "analysis": {"ok": True},
                            "structured": {"ok": True},
                        },
                    }
                ],
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        server.MODEL_CONTEXT_CALIBRATION_PATH.write_text(
            json.dumps(
                {
                    "qwen36a3b": {
                        "contextWindow": 32768,
                        "maxInputChars": 43000,
                        "structuredEditChars": 39000,
                        "measuredAt": "2026-05-23T21:00:00+08:00",
                    }
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        return subprocess.CompletedProcess(args=["measure"], returncode=0, stdout="ok", stderr="")

    try:
        server.MODEL_CONTEXT_CALIBRATION_PATH = temp_dir / "model-context-calibration.json"
        server.MODEL_CONTEXT_SELECTIONS_PATH = temp_dir / "model-contexts.json"
        server.LOGS_DIR = temp_dir / "logs"
        server.run_python_script_via_log = fake_runner
        result = server.run_model_context_calibration("qwen36a3b", [131072])
        assert_true(captured["script"] == "measure_context_limits.py", "calibration should run the benchmark script")
        assert_true(
            captured["args"] == [
                "--models",
                "qwen36a3b",
                "--contexts",
                "131072",
                "65536",
                "32768",
                "16384",
                "8192",
                "4096",
                "--stop-after-first-success",
            ],
            "calibration should test the upper context limit downward and stop after the first success",
        )
        assert_true(captured["timeout"] == 3600, "calibration should allow long-running model probes")
        assert_true(result["structuredEditChars"] == 39000, "calibration result should expose measured edit budget")
        try:
            server.run_model_context_calibration("invalid", [32768])
            raise AssertionError("invalid model aliases should not fall back to the default model")
        except ValueError:
            pass

        def failing_runner(script_name, *args, timeout_seconds=None):
            server.LOGS_DIR.mkdir(parents=True, exist_ok=True)
            (server.LOGS_DIR / "model-context-bench.json").write_text(
                json.dumps(
                    [
                        {
                            "model": "qwen36a3b",
                            "context": 32768,
                            "startup_ok": False,
                            "startup_output": "failed to allocate buffer for kv cache",
                            "tests": {},
                        }
                    ],
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )
            return subprocess.CompletedProcess(args=["measure"], returncode=0, stdout="ok", stderr="")

        server.run_python_script_via_log = failing_runner
        try:
            server.run_model_context_calibration("qwen36a3b", [32768])
            raise AssertionError("calibration with no successful context should fail instead of showing empty results")
        except RuntimeError as exc:
            assert_true("No context calibration succeeded" in str(exc), "empty calibration should explain that no context succeeded")
    finally:
        server.MODEL_CONTEXT_CALIBRATION_PATH = old_calibration_path
        server.MODEL_CONTEXT_SELECTIONS_PATH = old_selection_path
        server.LOGS_DIR = old_logs_dir
        server.run_python_script_via_log = old_runner
        shutil.rmtree(temp_dir, ignore_errors=True)


def test_context_benchmark_reuses_existing_compatible_model_server():
    old_request_json = measure_context_limits.request_json
    old_request_chat = measure_context_limits.request_chat
    old_kill = measure_context_limits.kill_bench_ports
    old_run = measure_context_limits.subprocess.run
    calls = {"kill": 0, "run": 0, "chat": 0}

    def fake_request_json(port, path, timeout=30):
        if port == server.get_model_port("qwen36a3b") and path == "/v1/models":
            return {"ok": True, "data": {"data": [{"id": "qwen36a3b-local", "aliases": ["qwen36a3b-local"]}]}}
        if port == server.get_model_port("qwen36a3b") and path == "/props":
            return {"ok": True, "data": {"default_generation_settings": {"n_ctx": 131072}}}
        return {"ok": False, "error": "unexpected request"}

    def fake_request_chat(port, model, messages, max_tokens, timeout=180):
        calls["chat"] += 1
        return {"ok": True, "reply": "ok", "finish_reason": "stop", "length": 2}

    def fake_kill():
        calls["kill"] += 1

    def fake_run(*args, **kwargs):
        calls["run"] += 1
        raise AssertionError("probe should reuse existing compatible qwen36a3b server instead of launching another copy")

    try:
        measure_context_limits.request_json = fake_request_json
        measure_context_limits.request_chat = fake_request_chat
        measure_context_limits.kill_bench_ports = fake_kill
        measure_context_limits.subprocess.run = fake_run
        result = measure_context_limits.probe_model("qwen36a3b", 131072, 18120)
        assert_true(result["startup_ok"] is True, "reused qwen36a3b server should count as startup ok")
        assert_true(result.get("reused_existing_server") is True, "probe should mark reused existing server")
        assert_true(result["port"] == server.get_model_port("qwen36a3b"), "probe should use the existing model server port")
        assert_true(calls["kill"] == 0 and calls["run"] == 0, "probe should not kill ports or start another large model")
        assert_true(calls["chat"] >= 3, "probe should still run benchmark prompts against the existing server")
    finally:
        measure_context_limits.request_json = old_request_json
        measure_context_limits.request_chat = old_request_chat
        measure_context_limits.kill_bench_ports = old_kill
        measure_context_limits.subprocess.run = old_run


def test_start_server_uses_low_memory_model_env_for_qwen36_on_low_vram():
    resolved = subprocess.run(
        [str(ROOT / "runtime" / "WinPython" / "python" / "python.exe"), str(ROOT / "scripts" / "resolve_model_env.py"), "qwen36a3b"],
        cwd=str(ROOT),
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        check=False,
        timeout=60,
    )
    assert_true(resolved.returncode == 0, "resolve_model_env.py should resolve qwen36a3b")
    assert_true('MODEL_N_GPU_LAYERS=0' in resolved.stdout, "qwen36a3b should avoid GPU offload on this low-VRAM profile")
    assert_true('--batch-size=256' in resolved.stdout and '--ubatch-size=64' in resolved.stdout, "qwen36a3b should use low-memory batch settings")
    assert_true('--mlock' not in resolved.stdout, "qwen36a3b low-memory launch should not mlock the model")
    start_server = (ROOT / "scripts" / "start-server.cmd").read_text(encoding="utf-8")
    assert_true("timeout /t 2 /nobreak" in start_server and "Start-Sleep -Seconds 2" not in start_server, "start-server wait loop should avoid PowerShell sleep under OOM pressure")


def test_gemma_manifest_uses_unsloth_with_mmproj():
    manifest = json.loads((ROOT / "config" / "bootstrap.manifest.json").read_text(encoding="utf-8"))
    gemma = manifest["models"]["gemma4"]
    assert_true(gemma["repo"] == "unsloth/gemma-4-26B-A4B-it-GGUF", "gemma4 must use the Unsloth GGUF repo")
    removed_repo_owner = "bart" + "owski"
    assert_true(removed_repo_owner not in json.dumps(gemma), "gemma4 manifest must not reference the removed GGUF repo owner")
    assert_true(gemma["defaultQuant"] == "UD-Q4_K_M", "gemma4 should default to Unsloth UD-Q4_K_M")
    assert_true(gemma["supportsImages"] is True, "gemma4 should expose image support when mmproj is configured")
    assert_true(gemma["mmprojPatterns"], "gemma4 must require an mmproj file")


def test_new_model_catalog_exposes_hardware_metadata():
    payload = server.get_models_payload()
    models = payload["models"]
    for key in ("qwen36a3b", "qwen3coder30b", "glm46", "qwen25coder14b", "deepseekcoderlite"):
        assert_true(key in models, f"{key} should be available in the model catalog")
        assert_true(models[key].get("tier") in {"low", "standard", "high", "extreme"}, f"{key} should expose a hardware tier")
        assert_true(models[key].get("estimatedModelSizeGb", 0) > 0, f"{key} should expose estimated model size")
        assert_true(models[key].get("runtimeBackend"), f"{key} should expose selected runtime backend")
        assert_true(isinstance(models[key].get("recommended"), bool), f"{key} should expose recommendation state")
    assert_true(payload.get("hardwareProfile", {}).get("profile"), "/api/models should expose a hardware profile")
    assert_true(payload.get("recommendedModelKey") in models, "/api/models should expose a model recommendation")
    assert_true("--jinja" in models["glm46"].get("llamaArgs", []), "glm46 should require llama.cpp jinja chat templates")


def test_hardware_profile_classification_and_recommendations():
    from core.hardware import HardwareInfo, choose_recommended_model_key, classify_hardware, recommend_model_settings

    standard = HardwareInfo(
        total_ram_gb=32.0,
        cpu_cores=6,
        cpu_threads=12,
        gpus=[{"name": "AMD Radeon 760M Graphics", "vendor": "amd", "vramGb": 0.5}],
        has_nvidia_smi=False,
        has_vulkan=True,
    )
    standard_profile = classify_hardware(standard)
    assert_true(standard_profile["profile"] == "standard", "32GB RAM with AMD iGPU should be standard")
    standard_settings = recommend_model_settings(standard_profile, server.get_model_manifest("qwen25coder14b"))
    assert_true(standard_settings["runtimeBackend"] == "vulkan", "AMD iGPU with Vulkan should prefer Vulkan")
    assert_true(standard_settings["contextWindow"] == 32768, "standard hardware should use qwen25coder14b default context")

    high = HardwareInfo(
        total_ram_gb=64.0,
        cpu_cores=12,
        cpu_threads=24,
        gpus=[{"name": "NVIDIA GeForce RTX 4090", "vendor": "nvidia", "vramGb": 24.0}],
        has_nvidia_smi=True,
        has_vulkan=True,
    )
    high_profile = classify_hardware(high)
    assert_true(high_profile["profile"] == "high", "64GB RAM with NVIDIA GPU should be high")
    high_settings = recommend_model_settings(high_profile, server.get_model_manifest("qwen3coder30b"))
    assert_true(high_settings["runtimeBackend"] == "cuda", "NVIDIA GPU should prefer CUDA")
    assert_true(high_settings["contextWindow"] >= 65536, "high hardware should keep a large coding context")

    extreme = HardwareInfo(
        total_ram_gb=192.0,
        cpu_cores=24,
        cpu_threads=48,
        gpus=[{"name": "NVIDIA RTX 6000 Ada", "vendor": "nvidia", "vramGb": 48.0}],
        has_nvidia_smi=True,
        has_vulkan=True,
    )
    assert_true(classify_hardware(extreme)["profile"] == "extreme", "192GB RAM or 48GB VRAM should be extreme")

    low = HardwareInfo(
        total_ram_gb=16.0,
        cpu_cores=4,
        cpu_threads=8,
        gpus=[],
        has_nvidia_smi=False,
        has_vulkan=False,
    )
    assert_true(classify_hardware(low)["profile"] == "low", "low RAM without GPU should be low")


def test_qwen36a3b_catalog_targets_8gb_nvidia_moe_offload():
    from core.hardware import HardwareInfo, choose_recommended_model_key, classify_hardware, recommend_model_settings

    manifest = json.loads((ROOT / "config" / "bootstrap.manifest.json").read_text(encoding="utf-8"))
    qwen36 = manifest["models"]["qwen36a3b"]
    assert_true(qwen36["repo"] == "unsloth/Qwen3.6-35B-A3B-GGUF", "qwen36a3b should use the Unsloth GGUF repo")
    assert_true(qwen36["defaultQuant"] == "UD-Q4_K_M", "qwen36a3b should use the 8GB-oriented UD-Q4_K_M quant")
    assert_true("Qwen3.6-35B-A3B" in qwen36["displayName"], "qwen36a3b should expose the Qwen3.6 35B A3B model")
    assert_true(qwen36["supportsImages"] is True, "qwen36a3b should expose multimodal support")
    assert_true(qwen36["mmprojPatterns"], "qwen36a3b should require the vision mmproj file")
    assert_true(qwen36["contextWindow"] == 32768, "qwen36a3b should default to the 32k 8GB launch profile")
    assert_true(qwen36["nGpuLayers"] == 999, "qwen36a3b manual launcher should offload non-MoE layers")
    for arg in ("--flash-attn", "--jinja", "--n-cpu-moe=999", "--batch-size=512", "--ubatch-size=128", "--mlock"):
        assert_true(arg in qwen36["llamaArgs"], f"qwen36a3b should pass {arg} to llama.cpp")

    profile = classify_hardware(HardwareInfo(
        total_ram_gb=64.0,
        cpu_cores=12,
        cpu_threads=20,
        gpus=[{"name": "NVIDIA GeForce RTX 3070", "vendor": "nvidia", "vramGb": 8.0}],
        has_nvidia_smi=True,
        has_vulkan=True,
    ))
    raw_configs = {key: server.get_model_manifest(key) for key in server.SUPPORTED_MODEL_KEYS}
    assert_true(choose_recommended_model_key(profile, raw_configs) == "qwen36a3b", "64GB RAM plus RTX 3070 8GB should recommend qwen36a3b")
    settings = recommend_model_settings(profile, qwen36)
    assert_true(settings["runtimeBackend"] == "cuda", "qwen36a3b should use CUDA on NVIDIA")
    assert_true(settings["contextWindow"] == 32768, "qwen36a3b should keep the 32k context on 8GB-class NVIDIA")
    assert_true(settings["nGpuLayers"] == 999, "qwen36a3b should still offload non-MoE layers to GPU")

    low_vram_profile = classify_hardware(HardwareInfo(
        total_ram_gb=31.0,
        cpu_cores=6,
        cpu_threads=12,
        gpus=[{"name": "AMD Radeon 760M Graphics", "vendor": "amd", "vramGb": 0.5}],
        has_nvidia_smi=False,
        has_vulkan=True,
    ))
    low_vram_settings = recommend_model_settings(low_vram_profile, qwen36)
    assert_true(low_vram_settings["nGpuLayers"] == 0, "qwen36a3b should avoid iGPU offload on sub-4GB VRAM")
    low_vram_args = server.get_model_llama_args("qwen36a3b", low_vram_profile)
    assert_true("--batch-size=256" in low_vram_args, "qwen36a3b low-memory args should reduce batch size")
    assert_true("--ubatch-size=64" in low_vram_args, "qwen36a3b low-memory args should reduce ubatch size")
    assert_true("--mlock" not in low_vram_args, "qwen36a3b low-memory args should avoid mlock")


def test_llama_launcher_accepts_auto_hardware_args():
    source = (ROOT / "scripts" / "launch_llama_server.py").read_text(encoding="utf-8")
    assert_true('parser.add_argument("--n-gpu-layers"' in source, "launcher should accept --n-gpu-layers")
    assert_true('parser.add_argument("--n-cpu-moe"' in source, "launcher should accept llama.cpp MoE CPU offload")
    assert_true('parser.add_argument("--batch-size"' in source, "launcher should accept tuned logical batch size")
    assert_true('parser.add_argument("--ubatch-size"' in source, "launcher should accept tuned physical batch size")
    assert_true('parser.add_argument("--mlock"' in source, "launcher should accept model RAM lock")
    assert_true('parser.add_argument("--flash-attn"' in source, "launcher should accept --flash-attn")
    assert_true('parser.add_argument("--jinja"' in source, "launcher should accept --jinja")
    assert_true('"--flash-attn", "on"' in source, "launcher should pass an explicit Flash Attention value for current llama.cpp")
    assert_true('"--n-cpu-moe"' in source, "launcher should forward --n-cpu-moe to llama-server")
    assert_true('"--batch-size"' in source, "launcher should forward --batch-size to llama-server")
    assert_true('"--ubatch-size"' in source, "launcher should forward --ubatch-size to llama-server")
    server_source = (ROOT / "webui" / "server.py").read_text(encoding="utf-8")
    assert_true("append_llama_manifest_args" in server_source, "web server should centrally forward manifest llamaArgs")
    assert_true('"--n-cpu-moe"' in server_source, "web server should allow manifest MoE CPU offload args")
    assert_true('"--n-gpu-layers",\n            "0"' not in source, "launcher must not hard-code CPU-only GPU layers")
    assert_true("[CODEWORKER_LAUNCH_METADATA]" in source, "launcher should write detailed launch metadata into llama-server logs")


def test_launch_webui_restarts_stale_codeworker_server():
    source = (ROOT / "scripts" / "launch-webui.cmd").read_text(encoding="utf-8")
    assert_true(":reclaim_webui_port" in source, "launch-webui should reclaim an already-running CodeWorker Web UI server")
    assert_true("webui\\server.py" in source, "launch-webui should only reclaim the CodeWorker webui/server.py process")
    assert_true("foreach ($ownerPid in $listeners)" in source, "launch-webui should avoid the read-only PowerShell $PID variable")
    assert_true("Stop-Process -Id $ownerPid -Force" in source, "launch-webui should stop the stale CodeWorker process before restarting")
    assert_true("Port %WEBUI_PORT% is occupied by another process" in source, "launch-webui should refuse unknown processes on the Web UI port")


def test_bootstrap_stops_codeworker_runtime_users_before_winpython_update():
    source = (ROOT / "scripts" / "bootstrap.ps1").read_text(encoding="utf-8")
    assert_true("function Stop-CodeWorkerProcessesUsingRuntime" in source, "bootstrap should define a runtime process cleanup helper")
    assert_true('$Name -eq "winPython"' in source, "bootstrap should run cleanup before reinstalling WinPython")
    assert_true("Stop-CodeWorkerProcessesUsingRuntime -RootDir $RootDir -RuntimeDir $targetDir" in source, "WinPython update should stop CodeWorker processes that lock runtime DLLs")
    assert_true("Stop-Process -Id $process.ProcessId -Force" in source, "bootstrap cleanup should force-stop matching CodeWorker runtime processes")


def test_model_download_progress_payload_reports_file_percent():
    payload = server.build_model_download_progress_payload(
        "gemma4.gguf",
        bytes_written=5 * 1024 * 1024,
        total_bytes=20 * 1024 * 1024,
        file_index=0,
        file_count=2,
        segment_start=12,
        segment_end=92,
    )
    assert_true(payload["progress"] == 32, "overall task progress should map the file percent into the download segment")
    assert_true(payload["download"]["percent"] == 25, "download payload should expose exact current-file percent")
    assert_true(payload["download"]["fileName"] == "gemma4.gguf", "download payload should expose current file name")
    assert_true(payload["download"]["fileIndex"] == 1 and payload["download"]["fileCount"] == 2, "download payload should expose file position")
    assert_true("25%" in payload["message"], "download message should include visible percent")
    assert_true("5.0 MB" in payload["message"] and "20.0 MB" in payload["message"], "download message should include downloaded and total size")

    unknown = server.build_model_download_progress_payload(
        "model.gguf",
        bytes_written=3 * 1024 * 1024,
        total_bytes=0,
        file_index=0,
        file_count=1,
        segment_start=12,
        segment_end=92,
    )
    assert_true(unknown["progress"] == 12, "unknown content-length should keep segment start progress")
    assert_true(unknown["download"]["percent"] is None, "unknown content-length should explicitly expose null percent")
    assert_true("已下載 3.0 MB" in unknown["message"], "unknown content-length message should still show downloaded bytes")

    source = (ROOT / "webui" / "server.py").read_text(encoding="utf-8")
    assert_true(
        "ensure_local_model_server(model_key, port=get_model_port(model_key), task_id=task_id)" in source,
        "open project model preparation should pass task_id so first-time downloads can report file progress",
    )
    assert_true(
        "get_model_mmproj_patterns(model_key)" in source and "resolved_filenames.append(mmproj_filename)" in source,
        "model download progress should include a missing mmproj file when the selected model requires one",
    )
    assert_true(
        "download_model_with_progress(task_id, model_key, force=False)" in source and "if not force:" in source,
        "open-project model preparation should only download missing model files while manual redownload can still force refresh",
    )
    assert_true(
        "get_model_minimum_size_bytes(model_key)" in source and "Model file appears incomplete" in source,
        "webui model preparation should reject obvious partial GGUF downloads instead of treating them as ready",
    )

    bootstrap_source = (ROOT / "scripts" / "bootstrap.ps1").read_text(encoding="utf-8")
    assert_true(
        "$partDestination = \"$Destination.part\"" in bootstrap_source and "Move-Item -LiteralPath $partDestination -Destination $Destination -Force" in bootstrap_source,
        "bootstrap should download to a .part file before replacing the final GGUF",
    )
    assert_true(
        "--progress-bar" in bootstrap_source or "[DOWNLOAD]" in bootstrap_source,
        "bootstrap should show visible download progress for large model files",
    )
    assert_true(
        "estimatedModelSizeGb" in bootstrap_source and "looks incomplete" in bootstrap_source,
        "bootstrap should reject partial model files based on the manifest estimated size",
    )


def test_partial_model_file_is_rejected_by_size_guard():
    minimum = server.get_model_minimum_size_bytes("qwen25coder14b")
    assert_true(minimum > 4 * 1024 ** 3, "qwen25coder14b should have a useful minimum size guard")
    root = ROOT / ".tmp" / f"regression-partial-model-{uuid.uuid4().hex}"
    root.mkdir(parents=True, exist_ok=True)
    try:
        partial = root / "qwen2.5-coder-14b-instruct-q4_k_m.gguf"
        partial.write_bytes(b"GGUF")
        try:
            server.validate_model_file(partial, minimum)
        except ValueError as exc:
            assert_true("appears incomplete" in str(exc), "partial GGUF should produce a clear incomplete-file error")
        else:
            raise AssertionError("partial GGUF should not pass model validation")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_hardware_optimization_log_entry_contains_diagnostics():
    entry = server.build_hardware_optimization_log_entry(
        "model_launch_plan",
        {"profile": "high", "totalRamGb": 64, "gpuVendors": ["nvidia"], "maxVramGb": 24},
        "qwen3coder30b",
        models={"qwen3coder30b": {"tier": "high", "runtimeBackend": "cuda", "recommended": True}},
        selected_model_key="qwen3coder30b",
        auto_settings={"runtimeBackend": "cuda", "contextWindow": 131072, "nGpuLayers": 999, "threads": 16},
        launch_args=["python", "scripts/launch_llama_server.py", "--n-gpu-layers", "999"],
        llama_server_path=ROOT / "runtime" / "llama.cpp" / "llama-server.exe",
        model_path=ROOT / "models" / "qwen3-coder.gguf",
        log_path=ROOT / "logs" / "llama-server-qwen3coder30b.log",
        err_path=ROOT / "logs" / "llama-server-qwen3coder30b.err.log",
        port=8083,
        already_running=False,
        details="test",
    )
    assert_true(entry["event"] == "model_launch_plan", "hardware optimization log should include event")
    assert_true(entry["hardwareProfile"]["profile"] == "high", "hardware optimization log should include hardware profile")
    assert_true(entry["recommendedModelKey"] == "qwen3coder30b", "hardware optimization log should include recommended model")
    assert_true(entry["autoSettings"]["runtimeBackend"] == "cuda", "hardware optimization log should include backend")
    assert_true(entry["launch"]["port"] == 8083, "hardware optimization log should include launch port")
    assert_true("--n-gpu-layers" in entry["launch"]["argv"], "hardware optimization log should include launch argv")


def test_model_file_matching_does_not_fallback_on_pattern_miss():
    root = ROOT / ".tmp" / "regression-model-match"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "gemma-4-26B-A4B-it-UD-Q4_K_XL.gguf").write_bytes(b"fake")
    try:
        from core.models import match_first_model_file

        assert_true(
            match_first_model_file(root, ["*mmproj-BF16*.gguf", "*mmproj-F16*.gguf"]) is None,
            "mmproj lookup must not fall back to the main GGUF file",
        )
    finally:
        shutil.rmtree(root, ignore_errors=True)


def write_fixture_file(root, relative_path, content="x"):
    target = root / relative_path
    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(content, encoding="utf-8")


def test_project_structure_classifies_multi_language_files():
    root = ROOT / ".tmp" / "regression-project-structure"
    shutil.rmtree(root, ignore_errors=True)
    try:
        fixtures = [
            "Program.cs", "Game.csproj", "Form1.cs", "Form1.Designer.cs", "AudioManager.cs",
            "Assets/Sounds/a.wav", "bin/Debug/app.config", "obj/Debug/App.g.cs",
            "Project1.dpr", "Unit1.pas", "Unit1.dfm", "Project1.dproj",
            "main.cpp", "CMakeLists.txt", "include/app.hpp", "src/app.cpp", "App.vcxproj", "App.vcxproj.filters",
            "Program.vb", "Form1.vb", "Form1.Designer.vb", "App.vbproj",
            "package.json", "src/main.ts", "src/App.tsx", "src/App.test.tsx",
            "pyproject.toml", "app.py", "tests/test_app.py",
            "build.gradle.kts", "src/main/java/com/example/App.java", "src/test/java/com/example/AppTest.java",
            "go.mod", "cmd/tool/main.go", "Cargo.toml", "src/main.rs", "composer.json", "Gemfile",
        ]
        for path in fixtures:
            write_fixture_file(root, path)
        files = server.collect_project_files(root)
        structure = server.build_project_structure(root, files)
        categories = structure["categories"]
        recommended = structure["recommendedPins"]
        assert_true("Program.cs" in categories["entrypoints"], "C# Program.cs should be classified as an entrypoint")
        assert_true("Project1.dpr" in categories["entrypoints"], "Delphi .dpr should be classified as an entrypoint")
        assert_true("main.cpp" in categories["entrypoints"], "C++ main.cpp should be classified as an entrypoint")
        assert_true("Program.vb" in categories["entrypoints"], "VB Program.vb should be classified as an entrypoint")
        assert_true("Game.csproj" in categories["projectConfigs"], "C# project should be a project config")
        assert_true("Project1.dproj" in categories["projectConfigs"], "Delphi project should be a project config")
        assert_true("App.vcxproj" in categories["projectConfigs"], "C++ vcxproj should be a project config")
        assert_true("App.vbproj" in categories["projectConfigs"], "VB project should be a project config")
        assert_true("Form1.Designer.cs" in categories["uiFiles"], "WinForms designer should be classified as UI")
        assert_true("Unit1.dfm" in categories["uiFiles"], "Delphi DFM should be classified as UI")
        assert_true("Assets/Sounds/a.wav" in categories["assetFiles"], "asset audio should be classified as a resource")
        assert_true("src/App.test.tsx" in categories["testFiles"], "JS/TS test should be detected")
        assert_true("tests/test_app.py" in categories["testFiles"], "Python test should be detected")
        assert_true("bin/Debug/app.config" in categories["ignoredBuildOutputs"], "bin output should be ignored")
        assert_true("obj/Debug/App.g.cs" in categories["ignoredBuildOutputs"], "obj output should be ignored")
        assert_true("bin/Debug/app.config" not in recommended, "recommended pins must not include build output")
        assert_true(any(path in recommended for path in ["Program.cs", "Project1.dpr", "main.cpp", "Program.vb"]), "recommended pins should include entrypoints")
        assert_true("未找到測試檔案" not in structure["summary"], "summary should list tests when they exist")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_gemma_multimodal_payload_and_fallback():
    data = "data:image/png;base64," + base64.b64encode(b"fake-png").decode("ascii")
    upload = server.save_uploaded_file("pic.png", "image/png", data)
    original_has_native = server.model_has_native_image_transport
    try:
        server.model_has_native_image_transport = lambda key: False
        content, meta = server.build_attachment_chat_content("", "describe image", [upload], "gemma4")
        messages = server.prepare_messages_for_model("gemma4-local", server.build_raw_messages("gemma4", content, ""))
    finally:
        server.model_has_native_image_transport = original_has_native
    assert_true(meta["nativeImages"] == 0, "gemma4 should not waste time on native image payload without mmproj")
    assert_true(isinstance(messages[0]["content"], str), "gemma4 without mmproj should receive text fallback payload")


def test_image_metadata_fallback_blocks_guessing():
    data = "data:image/png;base64," + base64.b64encode(b"fake-png").decode("ascii")
    upload = server.save_uploaded_file("pic.png", "image/png", data)
    content, meta = server.build_attachment_chat_content("", "describe image", [upload], "gemma4", force_text_fallback=True)
    assert_true(meta["nativeImages"] == 0, "forced fallback should not include native images")
    assert_true(isinstance(content, str), "forced image fallback should produce a text prompt")
    assert_true("不得描述" in content and "不要猜測" in content, "image metadata fallback must explicitly block visual hallucination")


def test_video_metadata_fallback_blocks_guessing():
    upload = {
        "id": "video-1",
        "kind": "video",
        "name": "generated_video.mp4",
        "mimeType": "video/mp4",
        "sizeBytes": 1234,
        "extractionStatus": "video-keyframes-unavailable:ffmpeg-not-found",
        "durationSeconds": 0,
        "keyframeCount": 0,
    }
    block = server.build_attachment_prompt_block([upload], "gemma4")
    assert_true("不得猜測影片畫面" in block, "video metadata-only fallback must block visual guessing")


def test_video_timestamp_selection_handles_short_videos():
    assert_true(server.choose_video_timestamps(None, 3) == [0.1], "unknown duration should still try an early frame")
    short = server.choose_video_timestamps(0.4, 3)
    assert_true(short and short[0] <= 0.35, "short video should use a timestamp inside the clip")
    long = server.choose_video_timestamps(10.0, 3)
    assert_true(len(long) == 3 and long[0] == 0.1 and long[1] == 5.0, "long video should sample beginning/middle/end")
    budget, mode = server.choose_video_keyframe_budget(45.0, 12)
    assert_true(budget == 12 and mode == "balanced", "45s video should use balanced sampled keyframes")
    detailed = server.choose_video_timestamps(45.0, budget)
    assert_true(len(detailed) == 12 and detailed[0] == 0.1, "balanced video should sample more than beginning/middle/end")


def test_media_assessment_exposes_local_limits():
    assessment = server.get_media_analysis_assessment()
    assert_true("recommendedMaxKeyframes" in assessment, "media assessment should expose keyframe budget")
    assert_true("speechToText" in assessment, "media assessment should expose STT backend status")


def test_transcribe_media_attachment_updates_text_preview():
    root = ROOT / ".tmp" / "regression-stt"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "audio.wav"
    source.write_bytes(b"fake wav")
    upload = {
        "id": "audio-1",
        "kind": "audio",
        "name": "audio.wav",
        "mimeType": "audio/wav",
        "path": str(source),
        "extractionStatus": "audio-transcript-unavailable",
        "durationSeconds": 0,
        "textPreview": "",
        "textBlocks": [],
    }
    originals = (
        server.ensure_ffmpeg_runtime,
        server.media_has_audio_stream,
        server.get_video_duration_seconds,
        server.extract_media_audio_to_wav,
        server.get_stt_backend_status,
        server.transcribe_wav_with_backend,
    )
    try:
        server.ensure_ffmpeg_runtime = lambda: ("ffmpeg", "ffprobe", "ready")
        server.media_has_audio_stream = lambda source_path, ffprobe: True
        server.get_video_duration_seconds = lambda source_path, ffprobe: 2.5
        server.extract_media_audio_to_wav = lambda source_path, upload_id, ffmpeg: (source, "ready")
        server.get_stt_backend_status = lambda: {"available": True, "backend": "test"}
        server.transcribe_wav_with_backend = lambda wav_path: ("hello transcript", "test-backend")
        server.transcribe_media_attachment(upload)
    finally:
        (
            server.ensure_ffmpeg_runtime,
            server.media_has_audio_stream,
            server.get_video_duration_seconds,
            server.extract_media_audio_to_wav,
            server.get_stt_backend_status,
            server.transcribe_wav_with_backend,
        ) = originals
        shutil.rmtree(root, ignore_errors=True)
    assert_true(upload["textPreview"] == "hello transcript", "STT transcript should become attachment text")
    assert_true(upload["extractionStatus"] == "audio-transcript-extracted:test-backend", "audio extraction status should record STT backend")
    block = server.build_attachment_prompt_block([upload], "gemma4")
    assert_true("hello transcript" in block, "STT transcript should be sent to the model")


def test_history_continuation_uses_previous_answer_tail():
    history = [
        {"role": "user", "content": "請說明架構"},
        {"role": "assistant", "content": "<think>hidden reasoning</think>\n\n第一段答案\n第二段答案"},
    ]
    message = server.build_history_continuation_message("請繼續", history)
    assert_true(message is not None, "continue request should produce a history continuation prompt")
    assert_true("第一段答案" in message and "hidden reasoning" not in message, "continuation should use visible answer text, not reasoning")


def test_chat_messages_include_recent_history():
    history = [
        {"role": "user", "content": "上一題：想更新遊戲速度要怎麼修改？"},
        {"role": "assistant", "content": "<think>internal</think>\n\n請修改 Form1.cs 的 timer.Interval。"},
    ]
    messages = server.build_raw_messages("gemma4", "那上一題的檔案是哪個？", "system prompt", history=history)
    roles = [item["role"] for item in messages]
    assert_true(roles == ["system", "user", "assistant", "user"], "chat messages should include recent user/assistant history before the current user message")
    combined = "\n".join(str(item["content"]) for item in messages)
    assert_true("上一題" in combined and "timer.Interval" in combined, "recent history content should be visible to the model")
    assert_true("internal" not in combined, "history should strip reasoning blocks")


def test_chat_messages_include_compressed_memory_summary():
    summary = "使用者目標 / 待辦:\n- 想調整遊戲速度\n\n已提到檔案 / 符號:\n- Form1.cs\n- gameTimer.Interval"
    messages = server.build_raw_messages(
        "gemma4",
        "上一題要改哪裡？",
        "system prompt",
        history=[],
        memory_summary=summary,
    )
    assert_true(messages[0]["role"] == "system", "compressed memory should be added to the system prompt")
    assert_true("COMPRESSED CONVERSATION MEMORY" in messages[0]["content"], "system prompt should include compressed memory heading")
    assert_true("gameTimer.Interval" in messages[0]["content"], "compressed memory should preserve important implementation references")


def test_compact_session_memory_keeps_ui_history_and_builds_summary():
    old_values = (
        list(server.STATE.history),
        server.STATE.memory_summary,
        server.STATE.memory_compacted_count,
    )
    try:
        with server.STATE_LOCK:
            server.STATE.history = []
            server.STATE.memory_summary = ""
            server.STATE.memory_compacted_count = 0
            for index in range(8):
                server.STATE.history.append({"role": "user", "content": f"第 {index} 題：想更新遊戲速度，請看 Form1.cs"})
                server.STATE.history.append({"role": "assistant", "content": f"第 {index} 答：修改 gameTimer.Interval 與 gameSpeed。"})
            original_len = len(server.STATE.history)
            server.compact_session_memory_locked("gemma4")
            assert_true(len(server.STATE.history) == original_len, "memory compaction should not remove visible UI history")
            assert_true(server.STATE.memory_compacted_count > 0, "memory compaction should record the compacted boundary")
            assert_true("Form1.cs" in server.STATE.memory_summary and "gameTimer.Interval" in server.STATE.memory_summary, "memory summary should preserve important file and symbol references")
    finally:
        with server.STATE_LOCK:
            server.STATE.history, server.STATE.memory_summary, server.STATE.memory_compacted_count = old_values


def test_length_continuation_drops_large_project_context():
    class FakeResponse:
        def __init__(self, lines):
            self.lines = [line.encode("utf-8") for line in lines]

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return False

        def __iter__(self):
            return iter(self.lines)

    requests = []
    original_urlopen = server.urllib.request.urlopen
    try:
        def fake_urlopen(request, timeout=0):
            payload = json.loads(request.data.decode("utf-8"))
            requests.append(payload)
            if len(requests) == 1:
                return FakeResponse([
                    'data: {"choices":[{"delta":{"content":"first part"},"finish_reason":null}]}\n',
                    'data: {"choices":[{"delta":{},"finish_reason":"length"}]}\n',
                    "data: [DONE]\n",
                ])
            return FakeResponse([
                'data: {"choices":[{"delta":{"content":" second part"},"finish_reason":"stop"}]}\n',
                "data: [DONE]\n",
            ])

        server.urllib.request.urlopen = fake_urlopen
        events = list(server.stream_local_model_events(
            "gemma4-local",
            [
                {"role": "system", "content": "system prompt"},
                {"role": "user", "content": "PROJECT RAG CONTEXT\n" + ("large context\n" * 200)},
            ],
            timeout_seconds=1,
            max_tokens=16,
            continue_on_length=1,
        ))
    finally:
        server.urllib.request.urlopen = original_urlopen
    assert_true(len(requests) == 2, "length continuation should call the model twice")
    second_messages = requests[1]["messages"]
    second_payload_text = json.dumps(second_messages, ensure_ascii=False)
    assert_true("PROJECT RAG CONTEXT" not in second_payload_text, "length continuation should not resend large project context")
    assert_true("first part" in second_payload_text, "length continuation should include the previous answer tail")
    assert_true(any(event.get("type") == "content" and event.get("text") == " second part" for event in events), "length continuation should stream continued content")


def test_partial_stream_reply_can_be_saved_for_continue():
    old_values = (
        list(server.STATE.history),
        server.STATE.memory_summary,
        server.STATE.memory_compacted_count,
        server.STATE.model_key,
        server.STATE.model_alias,
    )
    try:
        with server.STATE_LOCK:
            server.STATE.history = []
            server.STATE.memory_summary = ""
            server.STATE.memory_compacted_count = 0
            partial = server.build_stream_reply_text("reasoning", "已輸出的部分回答")
            server.append_chat_exchange_locked("gemma4", "請長篇說明", [], "PROJECT RAG CONTEXT", partial, assistant_kind="chat-partial")
            continuation = server.build_history_continuation_message("請繼續", server.STATE.history)
        assert_true(continuation is not None, "partial stream output should be available for manual continue")
        assert_true("已輸出的部分回答" in continuation, "continue prompt should include partial answer tail")
        assert_true("PROJECT RAG CONTEXT" not in continuation, "manual continue should not re-inject full project context")
    finally:
        with server.STATE_LOCK:
            (
                server.STATE.history,
                server.STATE.memory_summary,
                server.STATE.memory_compacted_count,
                server.STATE.model_key,
                server.STATE.model_alias,
            ) = old_values


def test_stream_reasoning_only_length_retries_for_final_answer():
    class FakeResponse:
        def __init__(self, lines):
            self.lines = [line.encode("utf-8") for line in lines]

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return False

        def __iter__(self):
            return iter(self.lines)

    calls = {"count": 0}
    original_urlopen = server.urllib.request.urlopen
    try:
        def fake_urlopen(_request, timeout=0):
            calls["count"] += 1
            if calls["count"] == 1:
                return FakeResponse([
                    'data: {"choices":[{"delta":{"reasoning_content":"thinking only"},"finish_reason":null}]}\n',
                    'data: {"choices":[{"delta":{},"finish_reason":"length"}]}\n',
                    "data: [DONE]\n",
                ])
            return FakeResponse([
                'data: {"choices":[{"delta":{"content":"final answer"},"finish_reason":"stop"}]}\n',
                "data: [DONE]\n",
            ])

        server.urllib.request.urlopen = fake_urlopen
        events = list(server.stream_local_model_events(
            "gemma4-local",
            [{"role": "user", "content": "hello"}],
            timeout_seconds=1,
            max_tokens=16,
            continue_on_length=1,
        ))
    finally:
        server.urllib.request.urlopen = original_urlopen
    assert_true(calls["count"] == 2, "reasoning-only length response should retry once")
    assert_true(any(event.get("type") == "continuation" and "最終答案" in event.get("text", "") for event in events), "retry should explain answer-only continuation")
    assert_true(any(event.get("type") == "content" and event.get("text") == "final answer" for event in events), "retry should stream final answer content")


def test_gemma_native_image_payload_with_mmproj():
    data = "data:image/png;base64," + base64.b64encode(b"fake-png").decode("ascii")
    upload = server.save_uploaded_file("pic.png", "image/png", data)
    original_has_native = server.model_has_native_image_transport
    try:
        server.model_has_native_image_transport = lambda key: key == "gemma4"
        content, meta = server.build_attachment_chat_content("", "describe image", [upload], "gemma4")
        messages = server.prepare_messages_for_model("gemma4-local", server.build_raw_messages("gemma4", content, ""))
    finally:
        server.model_has_native_image_transport = original_has_native
    assert_true(meta["nativeImages"] == 1, "gemma4 with mmproj should send native image payload")
    assert_true(isinstance(messages[0]["content"], list), "gemma4 native image payload should use OpenAI multimodal content parts")


def test_prepare_attachments_does_not_use_qwen_helper():
    data = "data:image/png;base64," + base64.b64encode(b"fake-png").decode("ascii")
    upload = server.save_uploaded_file("pic.png", "image/png", data)
    original_ensure = server.ensure_local_model_server
    try:
        def fail_if_called(*args, **kwargs):
            raise AssertionError("prepare_attachments_for_model must not start a secondary vision model")

        server.ensure_local_model_server = fail_if_called
        prepared = server.prepare_attachments_for_model("gemma4", [upload])
    finally:
        server.ensure_local_model_server = original_ensure
    assert_true(prepared[0]["id"] == upload["id"], "prepare should preserve original image attachment")
    helper_status = "vision" + "-helper"
    assert_true(helper_status not in str(prepared[0].get("extractionStatus", "")), "prepare should not mark secondary vision status")


def test_stream_attachment_fallback_for_native_model():
    data = "data:image/png;base64," + base64.b64encode(b"fake-png").decode("ascii")
    upload = server.save_uploaded_file("pic.png", "image/png", data)
    original_has_native = server.model_has_native_image_transport
    original_stream = server.stream_local_model_events
    calls = {"count": 0}

    def fake_stream(model_alias, messages, timeout_seconds=180, max_tokens=600, continue_on_length=0):
        calls["count"] += 1
        if calls["count"] == 1:
            raise RuntimeError("Failed to call local model endpoint: HTTP 400: missing mmproj for image_url")
        yield {"type": "content", "text": "fallback ok"}
        yield {"type": "finish", "finishReason": "stop"}

    try:
        server.model_has_native_image_transport = lambda key: True
        server.stream_local_model_events = fake_stream
        events = list(
            server.stream_local_model_with_attachment_fallback(
                "qwen35-local",
                "qwen35",
                "",
                "describe image",
                [upload],
                "",
                max_tokens=32,
                timeout_seconds=5,
                continue_on_length=0,
            )
        )
    finally:
        server.model_has_native_image_transport = original_has_native
        server.stream_local_model_events = original_stream
    assert_true(any(item.get("type") == "attachment_fallback" for item in events), "stream fallback event was not emitted")
    assert_true(any(item.get("text") == "fallback ok" for item in events), "stream fallback did not continue to content")



def test_http_error_body_is_preserved():
    body = io.BytesIO(b'{"error":{"message":"missing mmproj for image_url"}}')
    exc = urllib.error.HTTPError("http://127.0.0.1", 400, "Bad Request", {}, body)
    try:
        server.raise_local_model_http_error(exc)
    except RuntimeError as err:
        assert_true("mmproj" in str(err) and "image_url" in str(err), "HTTPError body must be preserved for fallback detection")
        assert_true(server.is_multimodal_transport_error(err), "preserved HTTPError body should be detected as multimodal error")
    else:
        raise AssertionError("raise_local_model_http_error did not raise")


def test_rag_manifest_search_and_stale():
    root = ROOT / ".tmp" / "regression-project"
    data_dir = ROOT / ".tmp" / "regression-index"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    target = root / "main.py"
    target.write_text("import os\n\ndef target_login_flow(user):\n    return user.name\n", encoding="utf-8")
    (root / "notes.md").write_text("# Login\nThe target_login_flow lives in main.py.\n", encoding="utf-8")
    try:
        result = rebuild_index(root, data_dir)
        assert_true(result["files"] == 2, "RAG should index source and docs")
        manifest_path = Path(result["indexDir"]) / "manifest.json"
        assert_true(manifest_path.exists(), "RAG manifest.json was not written")
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        assert_true("sha256" in manifest["files"][0], "RAG manifest must include sha256")
        assert_true(not index_is_stale(root, data_dir), "fresh index should not be stale")
        matches = search_index(root, data_dir, "target_login_flow", limit=3)["matches"]
        assert_true(matches, "RAG search should find target_login_flow")
        assert_true(matches[0]["path"] == "main.py", "RAG search should return the matching path")
        assert_true(matches[0]["lineStart"] >= 1 and matches[0]["lineEnd"] >= matches[0]["lineStart"], "RAG match should include line range")
        target.write_text("import os\n\ndef target_login_flow(user):\n    return user.email\n", encoding="utf-8")
        assert_true(index_is_stale(root, data_dir), "modified file should make index stale")
    finally:
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_code_graph_indexes_symbols_and_relationships():
    root = ROOT / ".tmp" / "regression-code-graph"
    data_dir = ROOT / ".tmp" / "regression-code-graph-index"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "auth.py").write_text(
        "def validate_user(user):\n"
        "    return bool(user)\n\n"
        "def login(user):\n"
        "    return validate_user(user)\n\n"
        "class AuthService:\n"
        "    def sign_in(self, user):\n"
        "        return login(user)\n",
        encoding="utf-8",
    )
    try:
        result = rebuild_index(root, data_dir)
        graph = result.get("codeGraph", {})
        assert_true(int(graph.get("nodes", 0)) >= 4, "code graph should index file, functions, class, and method nodes")
        assert_true(int(graph.get("edges", 0)) >= 3, "code graph should create contains/calls edges")
        status = code_graph_status(root, data_dir)
        assert_true(status["ready"] is True, "code graph status should report ready after rebuild")
        search = search_code_graph(root, data_dir, "login", limit=5)
        names = {item["name"] for item in search.get("nodes", [])}
        assert_true("login" in names, "code graph search should find function symbols by name")
        context, coverage = build_code_graph_context(root, data_dir, "login sign_in", limit=5)
        assert_true("CODE GRAPH CONTEXT" in context, "code graph context should produce a compact markdown block")
        assert_true("sign_in" in context and "login" in context, "code graph context should include related symbols")
        assert_true(int(coverage.get("edgesSent", 0)) > 0, "code graph context should include relationship edges")
    finally:
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_codeworker_codegraph_plugin_is_installable_and_queryable():
    plugin_root = ROOT / "plugins" / "codeworker-codegraph"
    manifest = json.loads((plugin_root / ".codex-plugin" / "plugin.json").read_text(encoding="utf-8"))
    skill_path = plugin_root / "skills" / "codeworker-codegraph" / "SKILL.md"
    script_path = plugin_root / "scripts" / "query_codeworker_graph.py"
    skill_script_path = plugin_root / "skills" / "codeworker-codegraph" / "scripts" / "query_codeworker_graph.py"
    assert_true(manifest["name"] == "codeworker-codegraph", "plugin manifest should use the normalized plugin name")
    assert_true(manifest.get("skills") == "./skills/", "plugin manifest should expose the skills directory")
    skill_text = skill_path.read_text(encoding="utf-8")
    assert_true("name: codeworker-codegraph" in skill_text, "skill frontmatter should expose codeworker-codegraph")
    assert_true(script_path.exists(), "plugin should include the graph query helper script")
    assert_true(skill_script_path.exists(), "installed skill should be self-contained with its own query helper script")

    root = ROOT / ".tmp" / "regression-codegraph-plugin"
    data_dir = ROOT / ".tmp" / "regression-codegraph-plugin-index"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "sample.py").write_text("def plugin_entry():\n    return 1\n", encoding="utf-8")
    try:
        completed = subprocess.run(
            [
                sys.executable,
                str(script_path),
                "--project",
                str(root),
                "--data-dir",
                str(data_dir),
                "--rebuild",
                "--query",
                "plugin_entry",
                "--json",
            ],
            cwd=str(ROOT),
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=60,
            check=False,
        )
        assert_true(completed.returncode == 0, f"plugin query script should run successfully: {completed.stderr}")
        payload = json.loads(completed.stdout)
        assert_true(payload["coverage"]["nodesSent"] > 0, "plugin query script should return code graph context")
    finally:
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_codegraph_api_helpers_return_status_context_and_coverage():
    root = ROOT / ".tmp" / "regression-codegraph-api"
    data_dir = ROOT / ".tmp" / "regression-codegraph-api-index"
    old_data_dir = server.DATA_DIR
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "service.py").write_text(
        "def fetch_user(user_id):\n"
        "    return user_id\n\n"
        "def render_profile(user_id):\n"
        "    return fetch_user(user_id)\n",
        encoding="utf-8",
    )
    try:
        server.DATA_DIR = data_dir
        initial = server.code_graph_status_payload(root)
        assert_true(initial["ready"] is False, "CodeGraph status should report not ready before indexing")
        payload = server.code_graph_query_payload(root, "render_profile fetch_user", limit=5)
        assert_true(payload["status"]["ready"] is True, "CodeGraph query should ensure the graph index exists")
        assert_true(len(payload["nodes"]) >= 1, "CodeGraph query should return symbol nodes")
        assert_true("CODE GRAPH CONTEXT" in payload["context"], "CodeGraph query should return compact graph context")
        assert_true(int(payload["coverage"].get("nodesSent", 0)) >= 1, "CodeGraph query should return coverage nodes")
        assert_true("indexDir" in payload["status"], "CodeGraph status should include indexDir")
        assert_true(payload["matchedFiles"], "CodeGraph query should return matched files for pinning")
        assert_true(payload["message"], "CodeGraph query should return a user-facing message")
        status = server.code_graph_status_payload(root)
        assert_true(status["sampleSymbols"], "CodeGraph status should return sample symbols for query chips")
        assert_true(status["sampleFiles"], "CodeGraph status should return sample files for query chips")
        assert_true(status["indexUpdatedAt"], "CodeGraph status should return indexUpdatedAt")
        no_match = server.code_graph_query_payload(root, "definitely_missing_symbol_xyz", limit=5)
        assert_true(no_match["nodes"] == [] and no_match["edges"] == [], "missing symbol should not return false matches")
        assert_true(no_match["suggestions"], "missing symbol should return usable suggestions")
        assert_true("definitely_missing_symbol_xyz" in no_match["message"], "missing symbol message should include the query")
    finally:
        server.DATA_DIR = old_data_dir
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_rag_model_loading_locator_prefers_source_chunks():
    root = ROOT / ".tmp" / "regression-model-locator"
    data_dir = ROOT / ".tmp" / "regression-model-locator-index"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    (root / "webui").mkdir(parents=True, exist_ok=True)
    (root / "scripts").mkdir(parents=True, exist_ok=True)
    (root / "config").mkdir(parents=True, exist_ok=True)
    (root / "docs").mkdir(parents=True, exist_ok=True)
    (root / "runtime" / "WinPython" / "python" / "Lib").mkdir(parents=True, exist_ok=True)
    (root / "data" / "indexes" / "cached").mkdir(parents=True, exist_ok=True)
    (root / "webui" / "server.py").write_text(
        "def ensure_runtime_and_model(model_key):\n"
        "    model_file = resolve_model_file(model_key)\n"
        "    return model_file\n\n"
        "def ensure_local_model_server(model_key):\n"
        "    model_file = ensure_runtime_and_model(model_key)\n"
        "    return launch_llama_server(model_file)\n",
        encoding="utf-8",
    )
    (root / "scripts" / "launch_llama_server.py").write_text(
        "import subprocess\n\n"
        "def launch_llama_server(model_file, mmproj_file=None):\n"
        "    args = ['llama-server', '--model', str(model_file)]\n"
        "    if mmproj_file:\n"
        "        args += ['--mmproj', str(mmproj_file)]\n"
        "    return subprocess.Popen(args)\n",
        encoding="utf-8",
    )
    (root / "scripts" / "start-server.cmd").write_text(
        "@echo off\r\n"
        "set MODEL_FILE=%~1\r\n"
        "llama-server.exe --model \"%MODEL_FILE%\"\r\n",
        encoding="utf-8",
    )
    (root / "config" / "bootstrap.manifest.json").write_text(
        '{"models":{"gemma4":{"repo":"example/gemma","filePatterns":["*.gguf"]}}}',
        encoding="utf-8",
    )
    (root / "docs" / "model-notes.md").write_text(
        "# Model loading\nThis document mentions model loading but is not the implementation.\n",
        encoding="utf-8",
    )
    (root / "runtime" / "WinPython" / "python" / "Lib" / "noise.py").write_text(
        "def ensure_runtime_and_model():\n    return 'do not index bundled runtime'\n",
        encoding="utf-8",
    )
    (root / "data" / "indexes" / "cached" / "manifest.json").write_text(
        '{"summary":"do not index cached RAG output"}',
        encoding="utf-8",
    )
    try:
        result = rebuild_index(root, data_dir)
        assert_true(result["files"] == 5, "RAG should index project source, scripts, config, and docs only")
        matches = search_index(root, data_dir, "請問加載model的code在哪個檔案的哪一段？", limit=5)["matches"]
        assert_true(matches, "model loading locator query should return matches")
        assert_true(
            matches[0]["path"] in {"webui/server.py", "scripts/launch_llama_server.py", "scripts/start-server.cmd"},
            "model loading locator should prefer source code chunks over summaries",
        )
        assert_true(
            "llama-server" in matches[0]["content"] or "ensure_runtime_and_model" in matches[0]["content"],
            "top model loading match should include implementation content",
        )
        assert_true(
            all(not str(item["path"]).startswith(("runtime/", "data/indexes/")) for item in matches),
            "RAG search must not return bundled runtime or cached index files",
        )
    finally:
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_rag_chinese_game_speed_query_finds_code():
    root = ROOT / ".tmp" / "regression-game-speed"
    data_dir = ROOT / ".tmp" / "regression-game-speed-index"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "README.md").write_text(
        "# Game\n主要遊戲邏輯在 Form1.cs，速度可調整。\n",
        encoding="utf-8",
    )
    (root / "Form1.cs").write_text(
        "using System.Windows.Forms;\n\n"
        "public partial class Form1 : Form {\n"
        "    private Timer gameTimer = new Timer();\n"
        "    private int gameSpeed = 120;\n\n"
        "    private void StartGame() {\n"
        "        gameTimer.Interval = gameSpeed;\n"
        "        gameTimer.Tick += GameLoop;\n"
        "    }\n\n"
        "    private void GameLoop(object sender, System.EventArgs e) {\n"
        "        UpdatePlayer();\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    try:
        rebuild_index(root, data_dir)
        matches = search_index(root, data_dir, "想更新遊戲速度要怎麼修改？", limit=3)["matches"]
        assert_true(matches, "Chinese game-speed query should return matches")
        assert_true(matches[0]["path"] == "Form1.cs", "game-speed query should prefer implementation code over README")
        assert_true("gameTimer.Interval" in matches[0]["content"] or "gameSpeed" in matches[0]["content"], "top match should include speed implementation")
    finally:
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_project_rag_context_without_pins():
    root = ROOT / ".tmp" / "regression-chat-project"
    data_dir = ROOT / ".tmp" / "regression-chat-index"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "main.py").write_text("def target_login_flow(user):\n    return user.name\n", encoding="utf-8")
    old_data_dir = server.DATA_DIR
    server.DATA_DIR = data_dir
    try:
        files = server.collect_project_files(root)
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4-local",
            files=files,
            summary=server.build_summary(root, files, [], []),
            ui_state="ready",
        )
        context, coverage = server.build_project_rag_context(root, state, "target_login_flow 在哪裡", "gemma4")
        assert_true(coverage["mode"] == "project-rag", "chat without pins should use project-rag context")
        assert_true("target_login_flow" in context, "project-rag context should include matching chunk")
    finally:
        server.DATA_DIR = old_data_dir
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_edit_context_sends_full_pinned_file_when_budget_allows_existing_effect_changes():
    root = ROOT / ".tmp" / f"regression-edit-full-pinned-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Form1.cs"
    padding = "\n".join(f"    private int Padding{i}() {{ return {i}; }}" for i in range(220))
    source.write_text(
        "using System.Drawing;\n"
        "public class Form1 {\n"
        "    private readonly List<ClearParticle> clearParticles = new();\n"
        f"{padding}\n"
        "    private void SpawnClearParticles(int row)\n"
        "    {\n"
        "        for (int x = 0; x < BoardWidth; x++)\n"
        "        {\n"
        "            Color color = pieceColors[board[row, x]];\n"
        "            for (int i = 0; i < 5; i++)\n"
        "            {\n"
        "                float speedX = (float)(random.NextDouble() * 6.0 - 3.0);\n"
        "                float speedY = (float)(random.NextDouble() * -4.0 - 1.5);\n"
        "                clearParticles.Add(new ClearParticle(new PointF(x, row), new PointF(speedX, speedY), color, 24, 4));\n"
        "            }\n"
        "        }\n"
        "    }\n"
        "\n"
        "    private void UpdateClearParticles()\n"
        "    {\n"
        "        for (int i = clearParticles.Count - 1; i >= 0; i--)\n"
        "        {\n"
        "            ClearParticle particle = clearParticles[i];\n"
        "            particle.Velocity = new PointF(particle.Velocity.X * 0.98f, particle.Velocity.Y + 0.35f);\n"
        "        }\n"
        "    }\n"
        "\n"
        "    private sealed class ClearParticle {}\n"
        "}\n",
        encoding="utf-8",
    )
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="qwen36a3b",
            model_alias="qwen36a3b-local",
            files=[server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#")],
            tree=["Form1.cs"],
            pinned_files=["Form1.cs"],
            ui_state="ready",
        )
        context, allowed = server.build_edit_context(
            root,
            state,
            "請讓既有碎裂動畫外拋更遠、停留更久，並調整物理更新。",
        )
        assert_true(allowed == ["Form1.cs"], "pinned edit context should keep the pinned file as the candidate")
        assert_true("完整內容" in context, "small enough pinned files should be sent as full content")
        assert_true("private void SpawnClearParticles" in context, "full pinned context should include the particle spawn method")
        assert_true("float speedX" in context and "float speedY" in context, "full pinned context should include velocity setup")
        assert_true("private void UpdateClearParticles" in context, "full pinned context should include the physics update method")
        coverage = server.get_last_edit_context_coverage()
        assert_true(coverage["files"][0]["mode"] == "full", "coverage should record full file mode")
        assert_true(coverage["files"][0]["truncated"] is False, "coverage should mark full pinned file as untruncated")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_chat_context_sends_full_line_numbered_pinned_file_for_all_models_when_budget_allows():
    root = ROOT / ".tmp" / f"regression-chat-full-pinned-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    old_calibration_path = server.MODEL_CONTEXT_CALIBRATION_PATH
    calibration_path = root / "model-context-calibration.json"
    try:
        body = "\n".join(
            [
                "using System;",
                "public class Form1",
                "{",
                *[f"    private int filler{i};" for i in range(1, 56)],
                "    private void SpawnClearParticles(int row)",
                "    {",
                "        float speedX = 1.25f;",
                "        float speedY = -2.5f;",
                "    }",
                "}",
            ]
        )
        (root / "Form1.cs").write_text(body, encoding="utf-8")
        calibration_path.write_text(
            json.dumps(
                {
                    "qwen36a3b": {"contextWindow": 65536, "maxInputChars": 120000, "structuredEditChars": 100000, "measuredAt": "2026-05-24T00:00:00+08:00"},
                    "gemma4": {"contextWindow": 65536, "maxInputChars": 120000, "structuredEditChars": 100000, "measuredAt": "2026-05-24T00:00:00+08:00"},
                    "qwen35": {"contextWindow": 65536, "maxInputChars": 120000, "structuredEditChars": 100000, "measuredAt": "2026-05-24T00:00:00+08:00"},
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        server.MODEL_CONTEXT_CALIBRATION_PATH = calibration_path
        for model_key in ("qwen36a3b", "gemma4", "qwen35"):
            state = server.SessionState(
                project_path=str(root),
                model_key=model_key,
                model_alias=f"{model_key}-local",
                files=[server.ProjectFile(path="Form1.cs", size=len(body), language="C#")],
                pinned_files=["Form1.cs"],
                ui_state="ready",
            )
            context, coverage = server.build_project_context(root, state, "請精確列出 SpawnClearParticles 的行號與內容")
            assert_true("檔案: Form1.cs [完整內容]" in context, f"{model_key} pinned chat context should send full content")
            assert_true("private void SpawnClearParticles" in context and "float speedX" in context and "float speedY" in context, f"{model_key} full chat context should include the complete method")
            assert_true("   59:     private void SpawnClearParticles(int row)" in context, f"{model_key} full chat context should include 1-based line numbers")
            assert_true(coverage["fullCount"] == 1 and coverage["files"][0]["truncated"] is False, f"{model_key} coverage should mark the pinned file as full")
    finally:
        server.MODEL_CONTEXT_CALIBRATION_PATH = old_calibration_path
        shutil.rmtree(root, ignore_errors=True)


def test_chat_context_falls_back_to_excerpt_for_all_models_when_full_file_exceeds_budget():
    root = ROOT / ".tmp" / f"regression-chat-excerpt-pinned-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    old_calibration_path = server.MODEL_CONTEXT_CALIBRATION_PATH
    calibration_path = root / "model-context-calibration.json"
    try:
        body = "\n".join(
            [
                "using System;",
                "public class Form1",
                "{",
                *[f"    private string filler{i} = \"{i:04d}-" + ("x" * 90) + "\";" for i in range(1, 180)],
                "    private void SpawnClearParticles(int row)",
                "    {",
                "        float speedX = 1.25f;",
                "        float speedY = -2.5f;",
                "    }",
                "}",
            ]
        )
        (root / "Form1.cs").write_text(body, encoding="utf-8")
        calibration_path.write_text(
            json.dumps(
                {
                    "qwen36a3b": {"contextWindow": 4096, "maxInputChars": 4000, "structuredEditChars": 3500, "measuredAt": "2026-05-24T00:00:00+08:00"},
                    "gemma4": {"contextWindow": 4096, "maxInputChars": 4000, "structuredEditChars": 3500, "measuredAt": "2026-05-24T00:00:00+08:00"},
                    "qwen35": {"contextWindow": 4096, "maxInputChars": 4000, "structuredEditChars": 3500, "measuredAt": "2026-05-24T00:00:00+08:00"},
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        server.MODEL_CONTEXT_CALIBRATION_PATH = calibration_path
        for model_key in ("qwen36a3b", "gemma4", "qwen35"):
            state = server.SessionState(
                project_path=str(root),
                model_key=model_key,
                model_alias=f"{model_key}-local",
                files=[server.ProjectFile(path="Form1.cs", size=len(body), language="C#")],
                pinned_files=["Form1.cs"],
                ui_state="ready",
            )
            context, coverage = server.build_project_context(root, state, "請精確列出 SpawnClearParticles 的內容")
            assert_true("檔案: Form1.cs [節錄模式]" in context, f"{model_key} should fall back to excerpt mode when full file is over budget")
            assert_true("檔案: Form1.cs [完整內容]" not in context, f"{model_key} should not label oversized context as full")
            assert_true("private void SpawnClearParticles" in context and "float speedX" in context, f"{model_key} excerpt should still include the relevant member")
            assert_true(coverage["fullCount"] == 0 and coverage["excerptCount"] == 1, f"{model_key} coverage should record excerpt mode")
            assert_true(coverage["files"][0]["truncated"] is True, f"{model_key} coverage should mark oversized file as truncated")
    finally:
        server.MODEL_CONTEXT_CALIBRATION_PATH = old_calibration_path
        shutil.rmtree(root, ignore_errors=True)


def test_edit_context_includes_all_pinned_files_when_total_budget_allows():
    root = ROOT / ".tmp" / f"regression-edit-multi-pinned-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    first = root / "Form1.cs"
    second = root / "AudioManager.cs"
    first.write_text("public class Form1 { private void RenderBoard() { DrawBoard(); } }\n", encoding="utf-8")
    second.write_text("public class AudioManager { public void PlayLineClear() { } }\n", encoding="utf-8")
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=[
                server.ProjectFile(path="Form1.cs", size=first.stat().st_size, language="C#"),
                server.ProjectFile(path="AudioManager.cs", size=second.stat().st_size, language="C#"),
            ],
            tree=["Form1.cs", "AudioManager.cs"],
            pinned_files=["AudioManager.cs", "Form1.cs"],
            ui_state="ready",
        )
        context, allowed = server.build_edit_context(root, state, "請調整清行音效與畫面提示")
        assert_true(set(allowed) == {"AudioManager.cs", "Form1.cs"}, "all pinned files should remain candidates")
        assert_true("檔案: AudioManager.cs" in context and "檔案: Form1.cs" in context, "all pinned files should be sent")
        coverage = server.get_last_edit_context_coverage()
        assert_true(len(coverage["files"]) == 2, "coverage should list both pinned files")
        assert_true(all(item["mode"] == "full" for item in coverage["files"]), "small pinned files should use full mode")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_edit_candidate_resolution_uses_rag_and_codegraph_without_pins():
    root = ROOT / ".tmp" / f"regression-edit-rag-codegraph-{uuid.uuid4().hex}"
    data_dir = ROOT / ".tmp" / f"regression-edit-rag-codegraph-index-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "README.md").write_text("# Notes\n\nrendering notes\n", encoding="utf-8")
    source = root / "Form1.cs"
    source.write_text(
        "public class Form1 {\n"
        "    private void RenderParticles()\n"
        "    {\n"
        "        particleVelocity += gravity;\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    old_data_dir = server.DATA_DIR
    server.DATA_DIR = data_dir
    try:
        rebuild_index(root, data_dir)
        files = server.collect_project_files(root)
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=files,
            tree=[file.path for file in files],
            ui_state="ready",
        )
        context, allowed = server.build_edit_context(root, state, "調整 RenderParticles 的 gravity 影響")
        assert_true(allowed and allowed[0] == "Form1.cs", "RAG/CodeGraph should prefer the source file without pins")
        assert_true("RenderParticles" in context and "particleVelocity" in context, "edit context should include the matched member")
    finally:
        server.DATA_DIR = old_data_dir
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_project_rag_rebuilds_graphless_existing_index():
    root = ROOT / ".tmp" / "regression-graphless-index"
    data_dir = ROOT / ".tmp" / "regression-graphless-data"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(data_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "main.py").write_text("def graphless_entry():\n    return 1\n", encoding="utf-8")
    old_data_dir = server.DATA_DIR
    server.DATA_DIR = data_dir
    try:
        first = rebuild_index(root, data_dir)
        db_path = Path(first["database"])
        conn = sqlite3.connect(str(db_path))
        try:
            conn.execute("DELETE FROM code_edges")
            conn.execute("DELETE FROM code_nodes")
            conn.commit()
        finally:
            conn.close()
        result, rebuilt = server.ensure_project_index(root)
        assert_true(rebuilt is True, "existing RAG index without code graph should be rebuilt after upgrade")
        assert_true(int(result.get("codeGraph", {}).get("nodes", 0)) > 0, "rebuilt graphless index should include code graph nodes")
    finally:
        server.DATA_DIR = old_data_dir
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(data_dir, ignore_errors=True)


def test_generated_text_file_requires_confirmation():
    root = ROOT / ".tmp" / "regression-generate"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    old_project = server.STATE.project_path
    old_ui = server.STATE.ui_state
    try:
        with server.STATE_LOCK:
            server.STATE.project_path = str(root)
            server.STATE.ui_state = "ready"
        action = server.create_generated_file_preview(
            root,
            {
                "targetPath": "generated/sample.md",
                "title": "sample",
                "content": "# Sample\n\nhello",
            },
        )
        target = root / "generated" / "sample.md"
        assert_true(not target.exists(), "generated file must not be written before confirmation")
        server.confirm_generated_file(str(action["id"]))
        assert_true(target.exists(), "generated file should be written after confirmation")
        assert_true("hello" in target.read_text(encoding="utf-8"), "generated file content should match preview content")
    finally:
        with server.STATE_LOCK:
            server.STATE.project_path = old_project
            server.STATE.ui_state = old_ui
        shutil.rmtree(root, ignore_errors=True)


def test_edit_actions_apply_with_git_checkpoint_and_restore():
    root = ROOT / ".tmp" / f"regression-edit-actions-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    target = root / "main.py"
    target.write_text("def hello():\n    return 'old'\n", encoding="utf-8")
    try:
        init = server.ensure_project_git_repo(root)
        assert_true(init["ready"] is True, "edit apply should prepare a git repository")
        before = target.read_text(encoding="utf-8")
        diff = server.generate_diff("main.py", before, before.replace("'old'", "'new'"))
        action = server.create_edit_action(
            root,
            "patch_file",
            "main.py",
            summary="Change return value",
            diff=diff,
            operations=[{"search": "return 'old'", "replace": "return 'new'"}],
        )
        assert_true(action["status"] == "pending", "edit action should start pending")
        assert_true(target.read_text(encoding="utf-8") == before, "creating an edit action must not write the file")
        plan = {"mode": "precise", "request": "change hello", "summary": "Change return value", "actions": [action]}
        result = server.apply_edit_actions(root, plan, [str(action["id"])])
        assert_true("return 'new'" in target.read_text(encoding="utf-8"), "apply should modify the target file")
        assert_true(result["preEditCommit"] and result["postEditCommit"], "apply should create pre/post git checkpoints")
        assert_true("main.py" in result["changedFiles"], "apply should report changed files")
        restore = server.restore_git_checkpoint(root, str(result["preEditCommit"]))
        assert_true(restore["restored"] is True, "restore should return success")
        assert_true(target.read_text(encoding="utf-8") == before, "restore should revert to the pre-edit checkpoint")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_edit_action_security_rejects_unsafe_paths_and_stale_patches():
    root = ROOT / ".tmp" / f"regression-edit-security-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "main.py").write_text("alpha\n", encoding="utf-8")
    try:
        server.ensure_project_git_repo(root)
        for bad_path in ("../escape.py", ".git/config", "runtime/tool.py", "models/model.gguf", "data/indexes/cache.db"):
            try:
                server.create_edit_action(root, "create_file", bad_path, content="x")
            except ValueError:
                pass
            else:
                raise AssertionError(f"{bad_path} should be rejected")
        action = server.create_edit_action(
            root,
            "patch_file",
            "main.py",
            operations=[{"search": "missing", "replace": "beta"}],
            diff="",
        )
        try:
            server.apply_edit_actions(root, {"actions": [action]}, [str(action["id"])])
        except server.EditApplyError as exc:
            assert_true("search 片段必須剛好匹配 1 次" in str(exc), "stale patch should fail with a clear match error")
            assert_true(exc.result.get("preEditCommit"), "failed apply should expose the pre-edit checkpoint for restore")
            server.restore_git_checkpoint(root, str(exc.result["preEditCommit"]))
        else:
            raise AssertionError("stale patch should be rejected")
        assert_true((root / "main.py").read_text(encoding="utf-8") == "alpha\n", "failed patch must not modify the file")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_edit_action_supports_create_replace_delete_rename_and_command():
    root = ROOT / ".tmp" / f"regression-edit-kinds-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "replace.txt").write_text("old\n", encoding="utf-8")
    (root / "delete.txt").write_text("delete me\n", encoding="utf-8")
    (root / "rename.txt").write_text("rename me\n", encoding="utf-8")
    try:
        server.ensure_project_git_repo(root)
        actions = [
            server.create_edit_action(root, "create_file", "created.txt", content="created\n"),
            server.create_edit_action(root, "replace_file", "replace.txt", content="new\n"),
            server.create_edit_action(root, "delete_file", "delete.txt"),
            server.create_edit_action(root, "rename_file", "rename.txt", target_path="renamed.txt"),
            server.create_edit_action(root, "run_command", "", command="cmd /c echo command-ok"),
        ]
        result = server.apply_edit_actions(root, {"actions": actions}, [str(action["id"]) for action in actions])
        assert_true((root / "created.txt").read_text(encoding="utf-8") == "created\n", "create_file should write a new file")
        assert_true((root / "replace.txt").read_text(encoding="utf-8") == "new\n", "replace_file should overwrite after confirmation")
        assert_true(not (root / "delete.txt").exists(), "delete_file should remove the file")
        assert_true(not (root / "rename.txt").exists() and (root / "renamed.txt").exists(), "rename_file should move the file")
        command_result = next(item for item in result["appliedActions"] if item["kind"] == "run_command")
        assert_true(command_result["returncode"] == 0 and "command-ok" in command_result["stdout"], "run_command should capture stdout")
        server.restore_git_checkpoint(root, str(result["preEditCommit"]))
        assert_true(not (root / "created.txt").exists(), "restore should remove created files")
        assert_true((root / "replace.txt").read_text(encoding="utf-8") == "old\n", "restore should revert replacements")
        assert_true((root / "delete.txt").exists(), "restore should bring deleted files back")
        assert_true((root / "rename.txt").exists() and not (root / "renamed.txt").exists(), "restore should revert renames")
        command_only = server.create_edit_action(root, "run_command", "", command="cmd /c echo read-only-command")
        command_only_result = server.apply_edit_actions(root, {"actions": [command_only]}, [str(command_only["id"])])
        assert_true(command_only_result["preEditCommit"], "command-only apply should still create a pre-edit checkpoint")
        assert_true(not command_only_result["postEditCommit"], "command-only apply without file changes should not create an empty post checkpoint")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_edit_apply_returns_validation_command_suggestions():
    root = ROOT / ".tmp" / f"regression-edit-validation-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    (root / "Demo.csproj").write_text("<Project Sdk=\"Microsoft.NET.Sdk\"></Project>\n", encoding="utf-8")
    source = root / "Program.cs"
    source.write_text("class Program { static string Name() => \"old\"; }\n", encoding="utf-8")
    try:
        server.ensure_project_git_repo(root)
        action = server.create_edit_action(
            root,
            "patch_file",
            "Program.cs",
            operations=[{"search": 'static string Name() => "old";', "replace": 'static string Name() => "new";'}],
        )
        result = server.apply_edit_actions(root, {"actions": [action]}, [str(action["id"])])
        commands = result.get("validationCommands", [])
        assert_true(commands, "apply result should include validation command suggestions")
        assert_true(commands[0]["command"] == 'dotnet build "Demo.csproj"', "csproj edits should suggest dotnet build")
        assert_true(commands[0]["autoRun"] is False, "validation suggestions should not auto-run arbitrary build commands")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_advisory_edit_plan_keeps_local_context_when_model_patch_is_unsafe():
    root = ROOT / ".tmp" / f"regression-edit-advisory-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    readme = root / "README.md"
    readme.write_text("# Game\n\nKeyboard controls are documented here.\n", encoding="utf-8")
    source = root / "Form1.cs"
    source.write_text(
        "using System;\n"
        "public class Form1 {\n"
        "    private int[,] board = new int[20, 10];\n"
        "    private void ClearLines() {\n"
        "        for (int y = 0; y < 20; y++) {\n"
        "            if (board[y, 0] > 0) {\n"
        "                board[y, 0] = 0;\n"
        "            }\n"
        "        }\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    original_call_local_model = server.call_local_model
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="qwen35",
            model_alias="qwen35",
            files=[server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#")],
            tree=["Form1.cs"],
            ui_state="ready",
        )

        def fake_call_local_model(*_args, **_kwargs):
            return json.dumps({
                "summary": "實作方塊消除粒子效果",
                "needMoreContext": [],
                "suggestions": [
                    {
                        "path": "Form1.cs",
                        "target": "ClearLines",
                        "whyHere": "ClearLines 目前直接清除資料，需要先產生粒子效果。",
                        "before": "",
                        "after": "CreateExplosionEffect(x, y, GetColorFromId(board[y, x]));",
                        "notes": [],
                    }
                ],
            }, ensure_ascii=False)

        server.call_local_model = fake_call_local_model
        plan = server.create_advisory_edit_plan(
            root,
            state,
            "幫我將方塊消除的功能加上特效",
            ["Form1.cs"],
            failure_reason="patch JSON invalid",
        )
        suggestion = plan["suggestions"][0]
        assert_true(suggestion["verified"] is False, "unsafe advisory suggestions should be marked unverified")
        assert_true(suggestion["source"] == "model-unverified", "unsafe advisory suggestions should expose their source")
        assert_true("private void ClearLines" in suggestion["before"], "advisory fallback should show the real local source region")
        assert_true("CreateExplosionEffect" in suggestion["after"], "unsafe advisory output should remain visible for manual review")
        assert_true("未建立可直接套用" in "；".join(suggestion["notes"]), "unsafe advisory output should clearly say it is not directly applicable")
    finally:
        server.call_local_model = original_call_local_model
        shutil.rmtree(root, ignore_errors=True)


def test_precise_validation_failure_logs_raw_reply_and_returns_unverified_advisory():
    root = ROOT / ".tmp" / f"regression-edit-raw-log-{uuid.uuid4().hex}"
    log_dir = ROOT / ".tmp" / f"regression-edit-raw-log-output-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    shutil.rmtree(log_dir, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    log_dir.mkdir(parents=True, exist_ok=True)
    source = root / "Program.cs"
    source.write_text(
        "public static class Program\n"
        "{\n"
        "    public static string Greeting() => \"Hello\";\n"
        "}\n",
        encoding="utf-8",
    )
    original_call_local_model = server.call_local_model
    original_logs_dir = server.LOGS_DIR
    try:
        server.LOGS_DIR = log_dir
        state = server.SessionState(
            project_path=str(root),
            model_key="qwen35",
            model_alias="qwen35",
            files=[server.ProjectFile(path="Program.cs", size=source.stat().st_size, language="C#")],
            tree=["Program.cs"],
            pinned_files=["Program.cs"],
            ui_state="ready",
        )
        calls = {"count": 0}

        def fake_call_local_model(*_args, **_kwargs):
            calls["count"] += 1
            if calls["count"] == 1:
                return json.dumps({
                    "summary": "更新 Greeting",
                    "needMoreContext": [],
                    "edits": [
                        {
                            "path": "Program.cs",
                            "target": "Greeting",
                            "reason": "測試 unsafe search",
                            "notes": [],
                            "operations": [
                                {"search": "return \"Missing\";", "replace": "return \"Hi\";"}
                            ],
                        }
                    ],
                }, ensure_ascii=False)
            return json.dumps({
                "summary": "精準模式失敗後的文字建議",
                "needMoreContext": [],
                "suggestions": [
                    {
                        "path": "Program.cs",
                        "target": "Greeting",
                        "whyHere": "Greeting 產生文字",
                        "before": "return \"Missing\";",
                        "after": "return \"Hi\";",
                        "notes": [],
                    }
                ],
            }, ensure_ascii=False)

        server.call_local_model = fake_call_local_model
        plan = server.create_edit_plan(root, state, "把 Greeting 回傳文字改成 Hi")
        suggestion = plan["suggestions"][0]
        logs = list(log_dir.glob("edit-plan-raw-qwen35-*.log"))
        assert_true(logs, "precise validation failure should write the raw model reply")
        assert_true("return \\\"Missing\\\";" in logs[0].read_text(encoding="utf-8"), "raw log should contain the invalid search snippet")
        assert_true(plan["mode"] == "advisory", "unsafe precise patch should return advisory mode")
        assert_true(suggestion["verified"] is False, "advisory suggestion after validation failure should be unverified")
        assert_true("return \"Missing\";" in suggestion["missingSearchSnippet"], "suggestion should expose the missing search snippet")
    finally:
        server.call_local_model = original_call_local_model
        server.LOGS_DIR = original_logs_dir
        shutil.rmtree(root, ignore_errors=True)
        shutil.rmtree(log_dir, ignore_errors=True)


def test_model_precise_patch_creates_applyable_action_without_hardcoded_rule():
    root = ROOT / ".tmp" / f"regression-edit-model-patch-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Program.cs"
    source.write_text(
        "namespace Demo;\n\n"
        "public static class Program\n"
        "{\n"
        "    public static string Greeting()\n"
        "    {\n"
        "        return \"Hello\";\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    original_call_local_model = server.call_local_model
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="qwen35",
            model_alias="qwen35",
            files=[server.ProjectFile(path="Program.cs", size=source.stat().st_size, language="C#")],
            tree=["Program.cs"],
            ui_state="ready",
        )

        def fake_call_local_model(*_args, **_kwargs):
            return json.dumps({
                "summary": "更新 Greeting 回傳文字",
                "needMoreContext": [],
                "edits": [
                    {
                        "path": "Program.cs",
                        "target": "Greeting",
                        "reason": "依需求調整顯示文字",
                        "notes": [],
                        "operations": [
                            {
                                "search": "        return \"Hello\";",
                                "replace": "        return \"Hi\";",
                            }
                        ],
                    }
                ],
            }, ensure_ascii=False)

        server.call_local_model = fake_call_local_model
        plan = server.create_edit_plan(root, state, "把 Greeting 回傳文字改成 Hi")
        actions = plan.get("actions", [])
        assert_true(plan["mode"] == "precise", "model patch should produce a precise plan")
        assert_true(len(actions) == 1 and actions[0]["kind"] == "patch_file", "model patch should become a patch_file action")
        server.apply_single_edit_action(root, actions[0])
        assert_true('return "Hi";' in source.read_text(encoding="utf-8"), "model patch action should actually modify the file")
    finally:
        server.call_local_model = original_call_local_model
        shutil.rmtree(root, ignore_errors=True)


def test_malformed_model_patch_is_salvaged_when_search_replace_are_unique():
    root = ROOT / ".tmp" / f"regression-edit-model-salvage-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Program.cs"
    source.write_text(
        "namespace Demo;\n\n"
        "public static class Program\n"
        "{\n"
        "    public static string Greeting()\n"
        "    {\n"
        "        return \"Hello\";\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    original_call_local_model = server.call_local_model
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="qwen35",
            model_alias="qwen35",
            files=[server.ProjectFile(path="Program.cs", size=source.stat().st_size, language="C#")],
            tree=["Program.cs"],
            ui_state="ready",
        )

        def fake_call_local_model(*_args, **_kwargs):
            return (
                '{\n'
                '  "summary": "更新 Greeting 回傳文字",\n'
                '  "path": "Program.cs",\n'
                '  "target": "Greeting",\n'
                '  "reason": "依需求調整顯示文字",\n'
                '  "search": "        return \\"Hello\\";",\n'
                '  "replace": "        return \\"Hi\\";"\n'
            )

        server.call_local_model = fake_call_local_model
        plan = server.create_edit_plan(root, state, "把 Greeting 回傳文字改成 Hi")
        actions = plan.get("actions", [])
        assert_true(len(actions) == 1 and actions[0]["kind"] == "patch_file", "salvaged patch should become a patch_file action")
        server.apply_single_edit_action(root, actions[0])
        assert_true('return "Hi";' in source.read_text(encoding="utf-8"), "salvaged patch should actually modify the file")
    finally:
        server.call_local_model = original_call_local_model
        shutil.rmtree(root, ignore_errors=True)


def test_fallback_advisory_uses_pending_target_to_locate_real_method():
    root = ROOT / ".tmp" / f"regression-edit-fallback-target-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Form1.cs"
    source.write_text(
        "using System.Drawing;\n"
        "public class Form1 {\n"
        "    private readonly Color[] pieceColors =\n"
        "    {\n"
        "        Color.Transparent,\n"
        "        Color.Cyan,\n"
        "        Color.Blue,\n"
        "        Color.Red\n"
        "    };\n"
        "\n"
        "    private int ClearLines()\n"
        "    {\n"
        "        int cleared = 0;\n"
        "        for (int y = 19; y >= 0; y--)\n"
        "        {\n"
        "            cleared++;\n"
        "        }\n"
        "        return cleared;\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=[server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#")],
            tree=["Form1.cs"],
            ui_state="ready",
        )
        plan = server.build_fallback_advisory_plan(
            root,
            state,
            "幫我將方塊消除的功能加上特效",
            ["Form1.cs"],
            "Gemma 4 patch 回傳不合法 JSON",
            pending_edit={
                "summary": "定位 ClearLines",
                "edits": [{"path": "Form1.cs", "target": "ClearLines", "location": "", "beforeSnippet": ""}],
            },
            raw_reply="<think>cannot json</think>",
        )
        suggestion = plan["suggestions"][0]
        assert_true("private int ClearLines" in suggestion["before"], "fallback should use the pending target instead of the file header")
        assert_true("pieceColors" not in suggestion["before"], "fallback should not select the earlier color array when ClearLines is known")
        assert_true("約第" in suggestion["location"], "fallback should show the located line range")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_fallback_advisory_salvages_partial_json_without_noisy_failure_text():
    root = ROOT / ".tmp" / f"regression-edit-partial-json-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Form1.cs"
    source.write_text(
        "public class Form1 {\n"
        "    private int ClearLines()\n"
        "    {\n"
        "        int cleared = 0;\n"
        "        return cleared;\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=[server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#")],
            tree=["Form1.cs"],
            ui_state="ready",
        )
        raw_reply = (
            '```json\n{\n'
            '  "summary": "新增粒子效果",\n'
            '  "path": "Form1.cs",\n'
            '  "target": "ClearLines",\n'
            '  "whyHere": "ClearLines 是消除列的位置",\n'
            '  "before": "private int ClearLines()\\n    {\\n        int cleared = 0;\\n        return cleared;\\n    }",\n'
            '  "after": "private int ClearLines()\\n    {\\n        int cleared = 0;\\n        SpawnLineParticles();\\n'
        )
        plan = server.build_fallback_advisory_plan(
            root,
            state,
            "幫我將方塊消除的功能加上特效",
            ["Form1.cs"],
            r"EDIT_PLAN_SCHEMA_INVALID: 模型回傳不合法 JSON。原始回覆尾段已寫入 C:\tmp\gemma4-advisory.log",
            raw_reply=raw_reply,
        )
        suggestion = plan["suggestions"][0]
        assert_true("private int ClearLines" in suggestion["before"], "partial JSON salvage should keep before snippet")
        assert_true("SpawnLineParticles" in suggestion["after"], "partial JSON salvage should keep incomplete after snippet")
        assert_true("完整合法 JSON" in plan["failureReason"], "user-facing failure reason should be concise")
        assert_true("Expecting value" not in plan["failureReason"], "user-facing failure reason should not expose parser details")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_edit_plan_timeout_short_circuits_to_local_fallback():
    root = ROOT / ".tmp" / f"regression-edit-timeout-fallback-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Form1.cs"
    source.write_text(
        "public class Form1 {\n"
        "    private int ClearLines()\n"
        "    {\n"
        "        int cleared = 0;\n"
        "        return cleared;\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    original_precise = server.create_precise_edit_plan
    original_advisory = server.create_advisory_edit_plan
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=[server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#")],
            tree=["Form1.cs"],
            ui_state="ready",
        )

        def timeout_precise(*args, **kwargs):
            raise RuntimeError(f"本地模型回應已等到目前上限仍未完成。timeout={server.EDIT_PLAN_TIMEOUT_SECONDS}s。")

        def fail_if_called(*args, **kwargs):
            raise AssertionError("timeout fallback must not start a second advisory model call")

        server.create_precise_edit_plan = timeout_precise
        server.create_advisory_edit_plan = fail_if_called
        plan = server.create_edit_plan(
            root,
            state,
            "幫我將方塊消除的功能加上特效，讓方塊有碎成小塊掉落消失的感覺。",
        )
        suggestion = plan["suggestions"][0]
        assert_true(plan["mode"] == "advisory", "timeout should return an advisory fallback")
        assert_true(f"timeout={server.EDIT_PLAN_TIMEOUT_SECONDS}s" in plan["failureReason"], "timeout reason should remain visible")
        assert_true("ClearLines" in suggestion["target"], "fallback should still locate ClearLines")
        assert_true("private int ClearLines" in suggestion["before"], "fallback should include the local source region")
        assert_true("ClearParticle" in suggestion["after"], "fallback should include a concrete local effect scaffold")
    finally:
        server.create_precise_edit_plan = original_precise
        server.create_advisory_edit_plan = original_advisory
        shutil.rmtree(root, ignore_errors=True)


def test_timeout_fallback_can_create_applyable_tetris_clear_effect_patch():
    root = ROOT / ".tmp" / f"regression-edit-direct-tetris-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Form1.cs"
    source.write_text(
        "public class Form1 {\n"
        "        private const int BoardWidth = 10;\n"
        "        private const int BoardHeight = 20;\n"
        "        private const int CellSize = 30;\n"
        "        private readonly int[,] board = new int[BoardHeight, BoardWidth];\n"
        "        private readonly Color[] pieceColors = { Color.Transparent, Color.Red };\n"
        "        private readonly Random random = new();\n"
        "        private readonly System.Windows.Forms.Timer gameTimer = new();\n"
        "        private readonly System.Windows.Forms.Timer shakeTimer = new();\n"
        "        private readonly AudioManager audioManager = new(LandingSoundPath, LineClearSoundPath, BackgroundMusicPath);\n"
        "        private Point shakeOffset = Point.Empty;\n"
        "\n"
        "        public Form1()\n"
        "        {\n"
        "            gameTimer.Tick += (_, _) => GameTick();\n"
        "            shakeTimer.Interval = 16;\n"
        "            shakeTimer.Tick += ShakeTimer_Tick;\n"
        "            KeyDown += Form1_KeyDown;\n"
        "        }\n"
        "\n"
        "        private void StartNewGame()\n"
        "        {\n"
        "            shakeOffset = Point.Empty;\n"
        "            shakeTimer.Stop();\n"
        "\n"
        "            nextPiece = CreateRandomPiece();\n"
        "        }\n"
        "\n"
        "        private int ClearLines()\n"
        "        {\n"
        "            int cleared = 0;\n"
        "            for (int y = BoardHeight - 1; y >= 0; y--)\n"
        "            {\n"
        "                bool full = true;\n"
        "                if (!full)\n"
        "                {\n"
        "                    continue;\n"
        "                }\n"
        "\n"
        "                cleared++;\n"
        "                for (int row = y; row > 0; row--)\n"
        "                {\n"
        "                }\n"
        "            }\n"
        "\n"
        "            return cleared;\n"
        "        }\n"
        "\n"
        "        protected override void OnPaint(PaintEventArgs e)\n"
        "        {\n"
        "            DrawBoard(g);\n"
        "            DrawCurrentPiece(g);\n"
        "            DrawSidePanel(g);\n"
        "        }\n"
        "\n"
        "        private sealed class FallingPiece\n"
        "        {\n"
        "        }\n"
        "}\n",
        encoding="utf-8",
    )
    original_precise = server.create_precise_edit_plan
    original_advisory = server.create_advisory_edit_plan
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=[server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#")],
            tree=["Form1.cs"],
            ui_state="ready",
        )

        def fail_if_precise_called(*args, **kwargs):
            raise AssertionError("local direct fallback should not call the model-backed precise planner")

        def fail_if_called(*args, **kwargs):
            raise AssertionError("timeout fallback must not start a second advisory model call")

        server.create_precise_edit_plan = fail_if_precise_called
        server.create_advisory_edit_plan = fail_if_called
        plan = server.create_edit_plan(
            root,
            state,
            "幫我將方塊消除的功能加上特效，讓方塊有碎成小塊掉落消失的感覺。",
        )
        actions = plan.get("actions", [])
        assert_true(len(actions) == 1, "direct fallback should create one patch_file action")
        assert_true(actions[0]["kind"] == "patch_file", "direct fallback action should be patch_file")
        assert_true(plan.get("edits") and plan["edits"][0].get("beforeSnippet"), "direct fallback should include reviewable snippets")
        server.apply_single_edit_action(root, actions[0])
        updated = source.read_text(encoding="utf-8")
        assert_true("private readonly List<ClearParticle> clearParticles" in updated, "applied patch should add particle state")
        assert_true("SpawnClearParticles(y);" in updated, "applied patch should hook ClearLines before shifting rows")
        assert_true("DrawClearParticles(g);" in updated, "applied patch should render particles")
        assert_true("private sealed class ClearParticle" in updated, "applied patch should add particle type")
    finally:
        server.create_precise_edit_plan = original_precise
        server.create_advisory_edit_plan = original_advisory
        shutil.rmtree(root, ignore_errors=True)


def test_edit_plan_timeout_is_scaled_for_local_coding_models():
    assert_true(
        server.get_edit_plan_timeout_seconds("qwen25coder14b") >= 420,
        "Qwen2.5-Coder 14B edit planning should allow slow local generation to finish",
    )
    assert_true(
        server.get_edit_plan_timeout_seconds("qwen36a3b") >= 900,
        "Qwen3.6-35B-A3B edit planning should not use the short generic timeout",
    )
    assert_true(
        server.get_edit_plan_timeout_seconds("qwen3coder30b") >= 600,
        "larger coding models should get a longer edit planning timeout",
    )
    assert_true(
        server.get_edit_plan_timeout_seconds("qwen36a3b", context_chars=100_000, max_tokens=1200)
        > server.get_edit_plan_timeout_seconds("qwen36a3b"),
        "large edit contexts should extend the edit planning timeout",
    )
    assert_true(
        server.get_edit_plan_timeout_seconds("unknown-model") == server.EDIT_PLAN_TIMEOUT_SECONDS,
        "unknown models should keep the default edit planning timeout",
    )


def test_edit_plan_model_call_uses_streaming_and_scaled_timeout():
    root = ROOT / ".tmp" / f"regression-edit-streaming-timeout-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    source = root / "Program.cs"
    source.write_text(
        "namespace Demo;\n\n"
        "public static class Program\n"
        "{\n"
        "    public static string Greeting()\n"
        "    {\n"
        "        return \"Hello\";\n"
        "    }\n"
        "}\n",
        encoding="utf-8",
    )
    original_call_local_model = server.call_local_model
    calls = []
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="qwen36a3b",
            model_alias="qwen36a3b",
            files=[server.ProjectFile(path="Program.cs", size=source.stat().st_size, language="C#")],
            pinned_files=["Program.cs"],
            tree=["Program.cs"],
            ui_state="ready",
        )

        def fake_call_local_model(*_args, **kwargs):
            calls.append(kwargs)
            return json.dumps({
                "summary": "更新 Greeting 回傳文字",
                "needMoreContext": [],
                "edits": [
                    {
                        "path": "Program.cs",
                        "target": "Greeting",
                        "reason": "依需求調整顯示文字",
                        "notes": [],
                        "operations": [
                            {
                                "search": "        return \"Hello\";",
                                "replace": "        return \"Hi\";",
                            }
                        ],
                    }
                ],
            }, ensure_ascii=False)

        server.call_local_model = fake_call_local_model
        plan = server.create_edit_plan(root, state, "把 Greeting 回傳文字改成 Hi")
        assert_true(plan["mode"] == "precise", "streaming model call should still produce a precise plan")
        assert_true(calls and calls[0].get("stream") is True, "edit plan model calls should stream so partial replies can be observed")
        assert_true(
            calls[0].get("timeout_seconds", 0) >= server.get_edit_plan_timeout_seconds("qwen36a3b"),
            "edit plan model calls should use the qwen36a3b scaled timeout",
        )
    finally:
        server.call_local_model = original_call_local_model
        shutil.rmtree(root, ignore_errors=True)


def test_streaming_edit_plan_timeout_logs_partial_reply():
    original_stream = server.stream_local_model_events
    original_logs_dir = server.LOGS_DIR
    log_dir = ROOT / ".tmp" / f"regression-stream-timeout-log-{uuid.uuid4().hex}"
    shutil.rmtree(log_dir, ignore_errors=True)
    log_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.LOGS_DIR = log_dir

        def fake_stream(*_args, **_kwargs):
            yield {"type": "content", "text": '{"summary": "partial"'}
            raise RuntimeError("本地模型回應已等到目前上限仍未完成。timeout=900s。")

        server.stream_local_model_events = fake_stream
        try:
            server.call_local_model(
                "qwen36a3b",
                [{"role": "user", "content": "edit"}],
                timeout_seconds=900,
                max_tokens=900,
                stream=True,
            )
            raise AssertionError("streaming timeout with partial content should raise")
        except RuntimeError as exc:
            details = str(exc)
            assert_true("已收到部分模型回應" in details, "timeout reason should say partial content was captured")
            assert_true("edit-plan-partial" in details, "timeout reason should include the partial reply log path")
            logs = list(log_dir.glob("edit-plan-partial-qwen36a3b-*.log"))
            assert_true(logs and '{"summary": "partial"' in logs[0].read_text(encoding="utf-8"), "partial reply should be written to a log")
    finally:
        server.stream_local_model_events = original_stream
        server.LOGS_DIR = original_logs_dir
        shutil.rmtree(log_dir, ignore_errors=True)


def test_local_direct_edit_changes_down_key_to_soft_drop_and_ctrl_to_hard_drop():
    root = ROOT / ".tmp" / f"regression-edit-direct-keys-{uuid.uuid4().hex}"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    readme = root / "README.md"
    readme.write_text("# Game\n\nKeyboard controls are documented here.\n", encoding="utf-8")
    source = root / "Form1.cs"
    source.write_text(
        "public class Form1 {\n"
        "        private void Form1_KeyDown(object? sender, KeyEventArgs e)\n"
        "        {\n"
        "            if (isGameOver)\n"
        "            {\n"
        "                return;\n"
        "            }\n"
        "\n"
        "            switch (e.KeyCode)\n"
        "            {\n"
        "                case Keys.Left:\n"
        "                case Keys.A:\n"
        "                    MovePiece(-1, 0);\n"
        "                    break;\n"
        "                case Keys.Right:\n"
        "                case Keys.D:\n"
        "                    MovePiece(1, 0);\n"
        "                    break;\n"
        "                case Keys.Down:\n"
        "                case Keys.Control:\n"
        "                    HardDrop();\n"
        "                    break;\n"
        "                case Keys.S:\n"
        "                    if (MovePiece(0, 1))\n"
        "                    {\n"
        "                        score += 1;\n"
        "                    }\n"
        "\n"
        "                    break;\n"
        "                case Keys.Up:\n"
        "                    RotatePiece();\n"
        "                    break;\n"
        "            }\n"
        "        }\n"
        "}\n",
        encoding="utf-8",
    )
    original_precise = server.create_precise_edit_plan
    original_advisory = server.create_advisory_edit_plan
    original_build_context = server.build_edit_context
    try:
        state = server.SessionState(
            project_path=str(root),
            model_key="gemma4",
            model_alias="gemma4",
            files=[
                server.ProjectFile(path="README.md", size=readme.stat().st_size, language="Markdown"),
                server.ProjectFile(path="Form1.cs", size=source.stat().st_size, language="C#"),
            ],
            tree=["README.md", "Form1.cs"],
            ui_state="ready",
        )

        def fail_if_precise_called(*args, **kwargs):
            raise AssertionError("down key direct edit should not call the model-backed precise planner")

        def fail_if_advisory_called(*args, **kwargs):
            raise AssertionError("down key direct edit should not call advisory model fallback")

        server.create_precise_edit_plan = fail_if_precise_called
        server.create_advisory_edit_plan = fail_if_advisory_called
        server.build_edit_context = lambda *args, **kwargs: ("可編輯候選檔案:\nREADME.md", ["README.md"])
        plan = server.create_edit_plan(
            root,
            state,
            "幫我修改按方向鑑的下不要直接掉到底，改成按ctrl掉到底。",
        )
        assert_true(plan["mode"] == "precise", "down key request should create a precise direct edit plan")
        assert_true(len(plan.get("actions", [])) == 1, "down key request should create one patch action")
        assert_true(plan["actions"][0]["kind"] == "patch_file", "down key action should be patch_file")
        assert_true("Form1_KeyDown" in plan["edits"][0]["target"], "down key edit should target Form1_KeyDown")
        assert_true("case Keys.Down" in plan["edits"][0]["beforeSnippet"], "review snippet should include existing Down case")
        assert_true("case Keys.ControlKey" in plan["edits"][0]["afterSnippet"], "replacement should support ControlKey")
        server.apply_single_edit_action(root, plan["actions"][0])
        updated = source.read_text(encoding="utf-8")
        assert_true("case Keys.Down:\n                case Keys.S:" in updated, "Down should share soft-drop behavior with S")
        assert_true("case Keys.Control:\n                case Keys.ControlKey:" in updated, "Ctrl should own hard drop behavior")
        assert_true("case Keys.Down:\n                case Keys.Control:\n                    HardDrop();" not in updated, "Down must no longer trigger HardDrop")
        repeated_plan = server.create_edit_plan(
            root,
            state,
            "幫我修改按方向鑑的下不要直接掉到底，改成按ctrl掉到底。",
        )
        assert_true(not repeated_plan.get("actions"), "repeated down key request should not create another action after it is fixed")
        assert_true("已符合" in repeated_plan["summary"], "repeated down key request should report the project already matches")
    finally:
        server.create_precise_edit_plan = original_precise
        server.create_advisory_edit_plan = original_advisory
        server.build_edit_context = original_build_context
        shutil.rmtree(root, ignore_errors=True)


def test_generation_prompt_infers_multiple_documents_from_previous_answer():
    history = [
        {"role": "user", "content": "請說明功能流程與使用場景"},
        {"role": "assistant", "content": "<think>internal</think>\n\n功能流程：先分析，再產出。\n\n使用場景：報告與簡報。"},
    ]
    requests = server.parse_generation_requests(
        {"prompt": "把剛剛的說明與使用場景做成一個PPTX跟PDF檔"},
        history,
    )
    targets = {item["targetPath"] for item in requests}
    assert_true(any(target.endswith(".pptx") for target in targets), "PPTX request should create a .pptx preview")
    assert_true(any(target.endswith(".pdf") for target in targets), "PDF request should create a .pdf preview")
    assert_true(not any(target.endswith(".md") for target in targets), "multi-format document request must not fall back to .md")
    combined = "\n".join(str(item["content"]) for item in requests)
    assert_true("功能流程" in combined and "internal" not in combined, "generation should use visible previous assistant content")


def test_generation_prompt_infers_excel():
    requests = server.parse_generation_requests({"prompt": "把測試清單做成 Excel 試算表"})
    assert_true(len(requests) == 1, "Excel-only request should create one preview")
    assert_true(requests[0]["targetPath"].endswith(".xlsx"), "Excel request should create an .xlsx target")


def test_generation_word_prompt_uses_previous_answer():
    history = [
        {"role": "user", "content": "請說明 CodeWorker"},
        {"role": "assistant", "content": "<think>hidden</think>\n\nCodeWorker 是本機 AI 助理。"},
    ]
    requests = server.parse_generation_requests({"prompt": "幫我把說明生成word檔"}, history)
    assert_true(requests[0]["targetPath"].endswith(".docx"), "word request should create a .docx target")
    assert_true("CodeWorker 是本機 AI 助理" in requests[0]["content"], "word request should use previous assistant answer")
    assert_true("hidden" not in requests[0]["content"], "word generation should strip reasoning")


def test_generation_with_previous_keyword_is_not_continuation():
    prompt = "請把剛剛的回答生成word檔"
    assert_true(server.is_model_file_generation_request(prompt), "word export prompt should be detected as file generation")
    assert_true(not server.is_history_continuation_request(prompt), "word export prompt must not be treated as continuation")


def test_generation_common_text_aliases():
    cases = {
        "請把剛剛的回答生成txt檔": ".txt",
        "請把剛剛的回答生成純文字檔": ".txt",
        "請把剛剛的回答生成md檔": ".md",
        "請把剛剛的回答生成py檔": ".py",
        "請把剛剛的回答生成js檔": ".js",
        "請把剛剛的回答生成ts檔": ".ts",
        "請把剛剛的回答生成json檔": ".json",
        "請把剛剛的回答生成html檔": ".html",
        "請把剛剛的回答生成css檔": ".css",
        "請把剛剛的回答生成yaml檔": ".yaml",
        "請把剛剛的回答生成sql檔": ".sql",
        "請把剛剛的回答生成cs檔": ".cs",
    }
    history = [{"role": "assistant", "content": "# 測試內容\n\nhello"}]
    for prompt, extension in cases.items():
        assert_true(server.is_model_file_generation_request(prompt), f"{prompt} should be detected as file generation")
        requests = server.parse_generation_requests({"prompt": prompt}, history)
        assert_true(requests[0]["targetPath"].endswith(extension), f"{prompt} should create {extension}")


def test_generated_docx_and_text_previews_can_be_created():
    root = ROOT / ".tmp" / "regression-generate-docx-text"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    old_project = server.STATE.project_path
    old_ui = server.STATE.ui_state
    try:
        with server.STATE_LOCK:
            server.STATE.project_path = str(root)
            server.STATE.ui_state = "ready"
        docx_action = server.create_generated_file_preview(
            root,
            {
                "targetPath": "generated/sample.docx",
                "title": "sample",
                "content": "# Sample\n\nhello",
            },
        )
        text_action = server.create_generated_file_preview(
            root,
            {
                "targetPath": "generated/sample.py",
                "title": "sample",
                "content": "print('hello')\n",
            },
        )
        assert_true(Path(str(docx_action["tempPath"])).exists(), "docx preview should create a temporary docx")
        assert_true(not (root / "generated" / "sample.docx").exists(), "docx preview must not write before confirmation")
        assert_true(str(text_action["content"]).strip() == "print('hello')", "text preview should keep source content")
        assert_true(not (root / "generated" / "sample.py").exists(), "text preview must not write before confirmation")
        docx_result = server.confirm_generated_file(str(docx_action["id"]))
        text_result = server.confirm_generated_file(str(text_action["id"]))
        assert_true((root / "generated" / "sample.docx").exists(), "docx should be written after confirmation")
        assert_true((root / "generated" / "sample.py").read_text(encoding="utf-8").strip() == "print('hello')", "text file should be written after confirmation")
        assert_true(docx_result["exists"] is True and Path(str(docx_result["path"])).exists(), "docx confirm response should prove the file exists")
        assert_true(text_result["exists"] is True and Path(str(text_result["path"])).exists(), "text confirm response should prove the file exists")
        assert_true(int(docx_result["sizeBytes"]) > 0, "docx confirm response should include non-zero size")
        assert_true(str(docx_result["absoluteTargetPath"]).endswith("sample.docx"), "confirm response should include absolute target path")
    finally:
        with server.STATE_LOCK:
            server.STATE.project_path = old_project
            server.STATE.ui_state = old_ui
        shutil.rmtree(root, ignore_errors=True)


def test_inline_docx_generation_uses_pasted_content_without_model():
    prompt = (
        "請把上面的內容生成docx檔給我\n"
        "# CodeWorker 產品說明書\n\n"
        "## 1. 產品概述\n"
        "CodeWorker 是本機 AI 程式碼助理。\n"
    )
    requests = server.build_generation_requests_from_inline_prompt(prompt)
    assert_true(len(requests) == 1, "inline docx request should create one generation request")
    assert_true(requests[0]["targetPath"].endswith(".docx"), "inline docx request should create a .docx target")
    assert_true(requests[0]["title"] == "CodeWorker 產品說明書", "inline docx should use the pasted heading as title")
    assert_true("CodeWorker 是本機 AI 程式碼助理" in requests[0]["content"], "inline docx should use pasted content directly")


def test_previous_answer_docx_generation_uses_history_without_model():
    history = [
        {"role": "user", "content": "請寫產品說明"},
        {
            "role": "assistant",
            "content": "# CodeWorker 產品說明書\n\n## 1. 產品概述\nCodeWorker 是本機 AI 程式碼助理。",
        },
    ]
    prompt = "請把上面的內容生成docx檔給我"
    requests = server.build_generation_requests_without_model(prompt, history)
    assert_true(len(requests) == 1, "previous-answer docx request should create one direct generation request")
    assert_true(requests[0]["targetPath"].endswith(".docx"), "previous-answer docx request should create .docx")
    assert_true(requests[0]["title"] == "CodeWorker 產品說明書", "previous-answer docx should use assistant heading as title")
    assert_true("CodeWorker 是本機 AI 程式碼助理" in requests[0]["content"], "previous-answer docx should use assistant content directly")


def test_thread_continuation_generation_loads_requested_thread_history():
    old_threads_dir = server.THREADS_DIR
    old_active_thread_id = server.ACTIVE_THREAD_ID
    old_state = (
        server.STATE.project_path,
        server.STATE.model_key,
        server.STATE.model_alias,
        list(server.STATE.history),
        server.STATE.memory_summary,
        server.STATE.memory_compacted_count,
    )
    root = ROOT / ".tmp" / "regression-thread-generation"
    threads_dir = root / "threads"
    shutil.rmtree(root, ignore_errors=True)
    threads_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.THREADS_DIR = threads_dir
        server.ACTIVE_THREAD_ID = "other-thread"
        thread_history = [
            {"role": "user", "content": "請寫產品說明"},
            {"role": "assistant", "content": "# CodeWorker 產品說明書\n\nCodeWorker 是本機 AI 程式碼助理。"},
        ]
        server.save_thread_file(
            {
                "id": "target-thread",
                "title": "產品說明",
                "createdAt": 1,
                "updatedAt": 2,
                "modelKey": "gemma4",
                "modelName": "Gemma 4 26B",
                "projectPath": str(root),
                "history": thread_history,
                "memorySummary": "",
                "memoryCompactedCount": 0,
            }
        )
        with server.STATE_LOCK:
            server.STATE.history = []
            server.activate_thread_for_request_locked("target-thread")
            requests = server.build_generation_requests_without_model("請把上面的內容生成docx檔給我", list(server.STATE.history))
        assert_true(server.ACTIVE_THREAD_ID == "target-thread", "chat request should activate the requested thread")
        assert_true(len(requests) == 1, "generation should use history from the requested thread")
        assert_true("CodeWorker 是本機 AI 程式碼助理" in requests[0]["content"], "requested thread history should provide generation content")
    finally:
        server.THREADS_DIR = old_threads_dir
        server.ACTIVE_THREAD_ID = old_active_thread_id
        with server.STATE_LOCK:
            (
                server.STATE.project_path,
                server.STATE.model_key,
                server.STATE.model_alias,
                server.STATE.history,
                server.STATE.memory_summary,
                server.STATE.memory_compacted_count,
            ) = old_state
        shutil.rmtree(root, ignore_errors=True)


def test_thread_save_updates_default_title_and_restores_project_context():
    old_threads_dir = server.THREADS_DIR
    old_active_thread_id = server.ACTIVE_THREAD_ID
    old_state = (
        server.STATE.project_path,
        server.STATE.model_key,
        server.STATE.model_alias,
        server.STATE.summary,
        list(server.STATE.tree),
        list(server.STATE.files),
        list(server.STATE.entrypoints),
        list(server.STATE.tests),
        list(server.STATE.pinned_files),
        server.STATE.current_preview_path,
        list(server.STATE.history),
        list(server.STATE.transcript),
        server.STATE.memory_summary,
        server.STATE.memory_compacted_count,
        server.STATE.pending_edit,
        server.STATE.ui_state,
    )
    root = ROOT / ".tmp" / "regression-thread-state"
    threads_dir = root / "threads"
    shutil.rmtree(root, ignore_errors=True)
    threads_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.THREADS_DIR = threads_dir
        server.ACTIVE_THREAD_ID = "default-title-thread"
        server.save_thread_file(
            {
                "id": "default-title-thread",
                "title": "New chat",
                "createdAt": 1,
                "updatedAt": 1,
                "modelKey": "gemma4",
                "modelName": "Gemma 4 26B",
                "projectPath": "",
                "history": [],
                "memorySummary": "",
                "memoryCompactedCount": 0,
            }
        )
        with server.STATE_LOCK:
            server.STATE.project_path = str(root)
            server.STATE.model_key = "gemma4"
            server.STATE.model_alias = server.get_model_alias("gemma4")
            server.STATE.summary = "project summary"
            server.STATE.tree = ["src/app.py"]
            server.STATE.files = [server.ProjectFile(path="src/app.py", size=12, language="Python")]
            server.STATE.entrypoints = ["src/app.py"]
            server.STATE.tests = ["tests/test_app.py"]
            server.STATE.pinned_files = ["src/app.py"]
            server.STATE.current_preview_path = "src/app.py"
            server.STATE.history = [{"role": "user", "content": "請分析登入流程涉及哪些檔案？"}]
            server.STATE.transcript = []
            server.STATE.memory_summary = "memory"
            server.STATE.memory_compacted_count = 2
            server.STATE.pending_edit = {"id": "edit-1", "status": "preview"}
            server.STATE.ui_state = "ready"
            server.save_current_thread_locked()
        saved = server.load_thread_file(server.thread_path("default-title-thread")) or {}
        assert_true(saved.get("title") != "New chat", "default thread title should update after the first user message")
        assert_true(saved.get("projectPath") == str(root), "thread should persist project path")
        assert_true(saved.get("tree") == ["src/app.py"], "thread should persist file tree")
        assert_true(saved.get("pinnedFiles") == ["src/app.py"], "thread should persist pinned files")
        assert_true(saved.get("pendingEdit") == {"id": "edit-1", "status": "preview"}, "thread should persist pending edit preview")
        assert_true(saved.get("transcriptVersion") == 1, "thread should persist transcript version")
        assert_true(len(saved.get("transcript") or []) == 1, "thread should fallback to history when saving missing transcript")
        assert_true(saved["transcript"][0]["kind"] == "chat", "history fallback transcript item should be a chat item")

        with server.STATE_LOCK:
            server.STATE.project_path = "C:/leaked-project"
            server.STATE.summary = "leaked"
            server.STATE.tree = ["leaked.py"]
            server.STATE.files = [server.ProjectFile(path="leaked.py", size=1, language="Python")]
            server.STATE.pinned_files = ["leaked.py"]
            server.STATE.current_preview_path = "leaked.py"
            server.STATE.pending_edit = {"id": "leaked"}
            server.STATE.transcript = [{"id": "leaked", "kind": "tool-codegraph", "title": "leaked", "createdAt": 1, "data": {}}]
            server.apply_thread_to_state_locked(
                {
                    "id": "empty-thread",
                    "title": "Empty",
                    "modelKey": "gemma4",
                    "history": [],
                    "projectPath": "",
                    "memorySummary": "",
                    "memoryCompactedCount": 0,
                }
            )
            assert_true(server.STATE.project_path is None, "selecting a thread without projectPath should clear previous project path")
            assert_true(server.STATE.tree == [], "selecting a thread without tree should clear previous file tree")
            assert_true(server.STATE.pending_edit is None, "selecting a thread without pendingEdit should clear previous edit preview")
            assert_true(server.STATE.ui_state == "idle", "selecting a thread without project context should return UI to idle")
            assert_true(server.STATE.transcript == [], "selecting an empty legacy thread should clear previous transcript")
    finally:
        server.THREADS_DIR = old_threads_dir
        server.ACTIVE_THREAD_ID = old_active_thread_id
        with server.STATE_LOCK:
            (
                server.STATE.project_path,
                server.STATE.model_key,
                server.STATE.model_alias,
                server.STATE.summary,
                server.STATE.tree,
                server.STATE.files,
                server.STATE.entrypoints,
                server.STATE.tests,
                server.STATE.pinned_files,
                server.STATE.current_preview_path,
                server.STATE.history,
                server.STATE.transcript,
                server.STATE.memory_summary,
                server.STATE.memory_compacted_count,
                server.STATE.pending_edit,
                server.STATE.ui_state,
            ) = old_state
        shutil.rmtree(root, ignore_errors=True)


def test_missing_thread_project_path_is_cleared_instead_of_restored():
    old_state = (
        server.STATE.project_path,
        server.STATE.summary,
        list(server.STATE.tree),
        list(server.STATE.files),
        list(server.STATE.pinned_files),
        server.STATE.pending_edit,
        server.STATE.ui_state,
    )
    missing = ROOT / ".tmp" / f"missing-project-{uuid.uuid4().hex}"
    shutil.rmtree(missing, ignore_errors=True)
    try:
        with server.STATE_LOCK:
            server.apply_thread_to_state_locked(
                {
                    "id": "missing-path-thread",
                    "title": "Missing path",
                    "modelKey": "qwen25coder14b",
                    "projectPath": str(missing),
                    "summary": "stale summary",
                    "tree": ["Form1.cs"],
                    "files": [{"path": "Form1.cs", "size": 10, "language": "C#"}],
                    "pinnedFiles": ["Form1.cs"],
                    "pendingEdit": {"id": "stale"},
                    "uiState": "ready",
                    "history": [],
                }
            )
            assert_true(server.STATE.project_path is None, "missing project path should not be restored")
            assert_true(server.STATE.ui_state == "idle", "missing project path should return UI state to idle")
            assert_true(server.STATE.tree == [], "missing project path should clear stale file tree")
            assert_true(server.STATE.pinned_files == [], "missing project path should clear stale pinned files")
            assert_true(server.STATE.pending_edit is None, "missing project path should clear stale pending edit")
    finally:
        with server.STATE_LOCK:
            (
                server.STATE.project_path,
                server.STATE.summary,
                server.STATE.tree,
                server.STATE.files,
                server.STATE.pinned_files,
                server.STATE.pending_edit,
                server.STATE.ui_state,
            ) = old_state


def test_invalid_open_project_path_leaves_server_idle():
    old_state = (
        server.STATE.project_path,
        server.STATE.model_key,
        server.STATE.model_alias,
        server.STATE.summary,
        list(server.STATE.tree),
        list(server.STATE.files),
        server.STATE.ui_state,
    )
    missing = ROOT / ".tmp" / f"invalid-open-{uuid.uuid4().hex}"
    shutil.rmtree(missing, ignore_errors=True)
    task = server.create_task("open-project")
    try:
        with server.STATE_LOCK:
            server.STATE.project_path = "C:/stale-project"
            server.STATE.summary = "stale"
            server.STATE.tree = ["stale.cs"]
            server.STATE.files = [server.ProjectFile(path="stale.cs", size=1, language="C#")]
            server.STATE.ui_state = "ready"
        server.open_project_worker(task.id, str(missing), "qwen25coder14b")
        result = server.get_task(task.id)
        assert_true(result is not None and result.status == "failed", "missing path should fail the open-project task")
        assert_true(result.error and result.error.get("code") == "PROJECT_PATH_INVALID", "missing path should report PROJECT_PATH_INVALID")
        assert_true(server.STATE.project_path is None, "failed missing path should clear stale project path")
        assert_true(server.STATE.ui_state == "idle", "failed missing path should leave server idle, not stuck in error with stale path")
        assert_true(server.STATE.tree == [], "failed missing path should clear stale tree")
    finally:
        with server.STATE_LOCK:
            (
                server.STATE.project_path,
                server.STATE.model_key,
                server.STATE.model_alias,
                server.STATE.summary,
                server.STATE.tree,
                server.STATE.files,
                server.STATE.ui_state,
            ) = old_state


def test_cleanup_empty_threads_only_removes_e2e_threads():
    old_threads_dir = server.THREADS_DIR
    old_active_thread_id = server.ACTIVE_THREAD_ID
    old_history = list(server.STATE.history)
    old_transcript = list(server.STATE.transcript)
    old_memory = server.STATE.memory_summary
    old_compacted = server.STATE.memory_compacted_count
    old_pending = server.STATE.pending_edit
    root = ROOT / ".tmp" / "regression-thread-cleanup"
    threads_dir = root / "threads"
    shutil.rmtree(root, ignore_errors=True)
    threads_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.THREADS_DIR = threads_dir
        server.ACTIVE_THREAD_ID = "e2e-empty"
        server.save_thread_file({"id": "e2e-empty", "title": "UI 驗證 1", "history": [], "metadata": {"source": "webui-e2e"}, "modelKey": "gemma4", "createdAt": 1, "updatedAt": 1})
        server.save_thread_file({"id": "prefix-empty", "title": "UI 驗證 2", "history": [], "modelKey": "gemma4", "createdAt": 1, "updatedAt": 1})
        server.save_thread_file({"id": "user-empty", "title": "使用者空白", "history": [], "modelKey": "gemma4", "createdAt": 1, "updatedAt": 1})
        server.save_thread_file({"id": "e2e-with-history", "title": "UI 驗證 3", "history": [{"role": "user", "content": "keep"}], "metadata": {"source": "webui-e2e"}, "modelKey": "gemma4", "createdAt": 1, "updatedAt": 1})
        with server.STATE_LOCK:
            data = server.cleanup_empty_threads_locked()
        assert_true(data["deletedCount"] == 2, "cleanup should remove only empty E2E threads")
        assert_true(not server.thread_path("e2e-empty").exists(), "cleanup should delete metadata-marked empty E2E thread")
        assert_true(not server.thread_path("prefix-empty").exists(), "cleanup should delete UI verification empty thread")
        assert_true(server.thread_path("user-empty").exists(), "cleanup must not delete normal empty user thread")
        assert_true(server.thread_path("e2e-with-history").exists(), "cleanup must not delete E2E thread with history")
    finally:
        server.THREADS_DIR = old_threads_dir
        server.ACTIVE_THREAD_ID = old_active_thread_id
        server.STATE.history = old_history
        server.STATE.transcript = old_transcript
        server.STATE.memory_summary = old_memory
        server.STATE.memory_compacted_count = old_compacted
        server.STATE.pending_edit = old_pending
        shutil.rmtree(root, ignore_errors=True)


def test_tool_transcript_items_do_not_enter_model_history():
    old_threads_dir = server.THREADS_DIR
    old_active_thread_id = server.ACTIVE_THREAD_ID
    old_history = list(server.STATE.history)
    old_transcript = list(server.STATE.transcript)
    root = ROOT / ".tmp" / "regression-transcript-tools"
    threads_dir = root / "threads"
    shutil.rmtree(root, ignore_errors=True)
    threads_dir.mkdir(parents=True, exist_ok=True)
    try:
        server.THREADS_DIR = threads_dir
        server.ACTIVE_THREAD_ID = "transcript-thread"
        with server.STATE_LOCK:
            server.STATE.history = [{"role": "user", "content": "請說明 Form1"}]
            server.STATE.transcript = []
            item = server.append_tool_transcript_item_locked(
                "tool-codegraph",
                "CodeGraph 查詢",
                "找到 Form1.cs",
                {"operation": "query", "matchedFiles": ["Form1.cs"]},
            )
        saved = server.load_thread_file(server.thread_path("transcript-thread")) or {}
        assert_true(item["kind"] == "tool-codegraph", "tool transcript item should keep its kind")
        assert_true(len(saved.get("history") or []) == 1, "tool transcript item must not be appended to chat history")
        assert_true(len(saved.get("transcript") or []) == 1, "tool transcript item should be persisted in transcript")
        model_messages = server.build_history_messages(saved.get("history"), "gemma4")
        assert_true(all("找到 Form1.cs" not in str(msg.get("content", "")) for msg in model_messages), "tool transcript content must not enter model history")
    finally:
        server.THREADS_DIR = old_threads_dir
        server.ACTIVE_THREAD_ID = old_active_thread_id
        with server.STATE_LOCK:
            server.STATE.history = old_history
            server.STATE.transcript = old_transcript
        shutil.rmtree(root, ignore_errors=True)


def test_generic_previous_answer_file_generation_defaults_to_markdown():
    history = [
        {"role": "user", "content": "請寫產品說明"},
        {"role": "assistant", "content": "# CodeWorker 產品說明書\n\nCodeWorker 是本機 AI 程式碼助理。"},
    ]
    requests = server.build_generation_requests_without_model("請把上面的內容生成檔案給我", history)
    assert_true(len(requests) == 1, "generic file generation should still create a direct preview from previous answer")
    assert_true(requests[0]["targetPath"].endswith(".md"), "generic generated file should default to markdown")
    assert_true(not server.is_model_file_generation_request("請問檔案生成的 code 在哪裡？"), "source-code questions about file generation must not become generation actions")


def test_generation_without_project_uses_app_root_and_previous_answer():
    old_root_dir = server.ROOT_DIR
    old_state = (
        server.STATE.project_path,
        server.STATE.ui_state,
        list(server.STATE.history),
    )
    root = ROOT / ".tmp" / "regression-generate-no-project"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    try:
        server.ROOT_DIR = root
        with server.STATE_LOCK:
            server.STATE.project_path = None
            server.STATE.ui_state = "idle"
            server.STATE.history = [
                {"role": "user", "content": "請長篇說明"},
                {"role": "assistant", "content": "# 魔術方塊操作說明\n\n## Slide 1\n- 版本與 Hold 功能"},
            ]
            generation_root = server.get_generation_root_locked()
            requests = server.build_generation_requests_without_model("請將上面的內容生成doc、ppt、pdf給我", list(server.STATE.history))
        suffixes = {Path(str(item["targetPath"])).suffix.lower() for item in requests}
        assert_true(generation_root == server.ROOT_DIR.resolve(), "file generation without an open project should use the CodeWorker app root")
        assert_true(suffixes == {".docx", ".pptx", ".pdf"}, "previous-answer multi-format generation should create docx, pptx, and pdf requests")
        assert_true(all("版本與 Hold 功能" in str(item["content"]) for item in requests), "multi-format generation should use previous assistant content")
        actions = [server.create_generated_file_preview(generation_root, request) for request in requests]
        results = [server.confirm_generated_file(str(action["id"])) for action in actions]
        assert_true(all(Path(str(result["path"])).exists() for result in results), "confirm should write generated files even when no project is open")
        assert_true(all(str(result["path"]).startswith(str(root)) for result in results), "no-project generated files should stay under the app root")
        legacy_action = server.create_generated_file_preview(generation_root, requests[0])
        legacy_action.pop("rootPath", None)
        with server.GENERATED_FILE_ACTIONS_LOCK:
            server.GENERATED_FILE_ACTIONS[str(legacy_action["id"])] = legacy_action
        legacy_result = server.confirm_generated_file(str(legacy_action["id"]))
        assert_true(Path(str(legacy_result["path"])).exists(), "legacy previews without rootPath should still confirm using targetPath")
    finally:
        server.ROOT_DIR = old_root_dir
        with server.STATE_LOCK:
            server.STATE.project_path, server.STATE.ui_state, server.STATE.history = old_state
        shutil.rmtree(root, ignore_errors=True)


def test_generated_pdf_keeps_chinese_text_extractable():
    root = ROOT / ".tmp" / "regression-generate-pdf"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    target = root / "sample.pdf"
    try:
        server.write_pdf(target, "測試文件", "### 標題\n\n您好！我是 CodeWorker。\n\n- 支援 PDF\n- 支援 PPTX")
        from pypdf import PdfReader

        text = "\n".join(page.extract_text() or "" for page in PdfReader(str(target)).pages)
        assert_true("測試文件" in text and "您好" in text, "generated PDF should preserve extractable Chinese text")
        assert_true("ЁН" not in text, "generated PDF should not produce garbled CJK glyph text")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_document_generation_cleans_markdown_for_pptx():
    root = ROOT / ".tmp" / "regression-generate-pptx"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    target = root / "sample.pptx"
    try:
        server.write_pptx(target, "測試簡報", "### 核心功能\n\n- **本機模型**\n- `RAG` 搜尋")
        from pptx import Presentation

        texts = []
        for slide in Presentation(str(target)).slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        combined = "\n".join(texts)
        assert_true("核心功能" in combined and "本機模型" in combined, "PPTX should keep headings and bullets")
        assert_true("**" not in combined and "`" not in combined, "PPTX should not expose raw Markdown markers")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_document_generation_splits_long_pptx_sections():
    root = ROOT / ".tmp" / "regression-generate-pptx-long"
    shutil.rmtree(root, ignore_errors=True)
    root.mkdir(parents=True, exist_ok=True)
    target = root / "sample.pptx"
    try:
        long_content = "\n".join(f"- 第 {index} 項內容很長，用來確認投影片不會把全部文字塞在同一頁。" for index in range(1, 16))
        server.write_pptx(target, "長內容簡報", long_content)
        from pptx import Presentation

        presentation = Presentation(str(target))
        assert_true(len(presentation.slides) > 2, "long PPTX content should be split across multiple slides")
    finally:
        shutil.rmtree(root, ignore_errors=True)


def test_model_initiated_generation_uses_model_title_for_filename():
    prompt = "我要生成一個專案功能介紹的PPT文件"
    reply = "# CodeWorker 專案功能介紹\n\n- 本機模型服務\n- 全專案 RAG"
    requests = server.build_generation_requests_from_model_reply(prompt, reply)
    assert_true(len(requests) == 1, "model-initiated PPT request should create one request")
    assert_true(requests[0]["targetPath"].endswith(".pptx"), "PPT request should create a .pptx target")
    assert_true("CodeWorker-專案功能介紹" in requests[0]["targetPath"], "generated filename should come from the model title")
    assert_true("本機模型服務" in requests[0]["content"], "generated content should come from model reply")


def test_generation_system_prompt_is_only_added_for_generation_requests():
    normal_prompt = server.build_chat_system_prompt("gemma4")
    generation_prompt = server.build_chat_system_prompt("gemma4", file_generation_requested=True)
    assert_true("CodeWorker 會在你回答後建立檔案預覽" not in normal_prompt, "normal chat system prompt should not mention file preview")
    assert_true("CodeWorker 會在你回答後建立檔案預覽" in generation_prompt, "generation chat prompt should instruct model to prepare content")


def test_stream_chat_initializes_model_generation_flag():
    source = inspect.getsource(server.WebUIHandler.handle_chat_stream)
    assert_true("file_generation_requested = file_generation_intent" in source, "stream chat must initialize file_generation_requested before preview creation")
    assert_true("generation_root = get_generation_root_locked()" in source, "stream chat must choose a generation root before preview creation")
    assert_true(
        "build_chat_system_prompt(snapshot.model_key, file_generation_requested=file_generation_requested)" in source,
        "stream chat must pass the generation flag into the model system prompt",
    )


def test_static_ui_exposes_file_tree_layout_and_ai_busy_indicator():
    html = (ROOT / "webui" / "static" / "index.html").read_text(encoding="utf-8")
    css = (ROOT / "webui" / "static" / "styles.css").read_text(encoding="utf-8")
    js = "\n".join(path.read_text(encoding="utf-8") for path in sorted((ROOT / "webui" / "static" / "js").glob("app-*.js")))
    version = (ROOT / "VERSION").read_text(encoding="utf-8").strip()
    launch = (ROOT / "scripts" / "launch-webui.cmd").read_text(encoding="utf-8")
    assert_true(server.APP_VERSION == version, "server should read the Web UI version from VERSION")
    assert_true(server.get_status_payload()["appVersion"] == version, "/api/status should expose the running app version")
    assert_true("WEBUI_EXPECTED_VERSION" in launch and "/api/status" in launch and "?v=" in launch, "launch-webui should verify version and open a cache-busted URL")
    assert_true('id="sidebarStatusDetails"' in html, "secondary model and hardware details should be grouped in a disclosure")
    assert_true('id="fileTreeCount"' in html, "file tree should expose a result count")
    assert_true("fileMetaByPath" in js and "formatBytes(size)" in js, "file tree and pinned summary should show pinned file sizes")
    assert_true('id="aiActivity"' in html and 'id="chatBusyBar"' in html, "chat UI should include a visible busy indicator")
    assert_true('id="editPlanBtn"' in html and 'id="gitDiffBtn"' in html and 'id="gitCheckpointBtn"' in html, "chat UI should expose edit plan and Git safety actions")
    assert_true('id="contextCalibrateBtn"' in html and "測試此模型可送出 KB" in html, "chat UI should expose context capacity calibration")
    assert_true('id="projectControlDetails"' in html and 'class="project-control-summary"' in html, "project controls should be collapsible")
    assert_true('class="composer-meta-row"' in html, "composer should keep attachments in a compact metadata row")
    assert_true(html.find('id="chatInput"') < html.find('id="codeGraphToolbar"') < html.find('class="chat-footer-row"'), "CodeGraph toolbar should sit below the input and above primary actions")
    assert_true(html.find('class="chat-footer-row"') < html.find('id="contextWindowSelect"'), "Context selector should live in the bottom action row")
    assert_true('class="chat-input-wrap"' in html, "chat input should have a dedicated wrapper for its inline label/help")
    assert_true(html.find('class="chat-input-wrap"') < html.find('id="chatInputLabel"') < html.find('id="chatImagePreview"'), "chat input label/help should sit with the textarea instead of the bottom context controls")
    assert_true(".sidebar" in css and "flex-direction: column;" in css, "sidebar should use flex layout so hidden panels do not reserve extra row gaps")
    assert_true(".summary-panel" in css and "flex: 0 1 190px;" in css, "sidebar should reserve useful summary height without adding phantom spacing")
    assert_true(".ai-spinner" in css and "@keyframes aiBusyBar" in css, "busy indicator should have spinner/bar animation styles")
    assert_true(".project-control-summary" in css and ".control-panel[open]" in css, "CSS should style collapsible project controls with bounded open height")
    assert_true(".chat-input-wrap" in css and ".composer-input-label" in css, "CSS should place the chat input label near the textarea")
    assert_true("input, select, textarea, button { font: inherit; font-size: 0.9rem; }" in css, "general controls should use the smaller UI text size")
    assert_true("padding: 7px 10px;" in css and "padding: 8px 10px;" in css, "smaller text should be paired with smaller button and input boxes")
    assert_true(".chat-content { white-space: pre-wrap; font-size: 0.9rem;" in css, "chat message text should use the smaller transcript size")
    assert_true(".diff-block" in css and ".edit-action-card" in css and ".edit-snippet-block" in css, "edit plan UI should style diff, snippets, and action cards")
    assert_true("function setAiBusy" in js, "app should control the AI busy indicator from JS")
    assert_true("setAiBusy(true" in js and "setAiBusy(false" in js, "chat/analyze flows should toggle the AI busy indicator")
    assert_true("AI 正在產生修改建議" in js and "finally" in js, "edit plan generation should show and clear the AI busy indicator")
    assert_true("function renderEditDetailHtml" in js and "建議替換前片段" in js and "建議替換後片段" in js, "edit plan UI should show location and before/after snippets")
    assert_true("未驗證參考片段" in js and "model-unverified" in js, "advisory edit UI should mark unverified model suggestions")
    assert_true("data.plan?.contextCoverage" in js and "修改計畫上下文" in js, "edit plan UI should show context coverage and full/region/window details")
    assert_true("body: JSON.stringify({ message, modelKey })" in js, "edit plan UI should send the currently selected model key")
    assert_true("/api/models/context-calibration" in js and "structuredEditChars" in js, "UI should expose measured structured edit budget")
    assert_true("function applyEditPlan" in js and "/api/edit/apply" in js, "UI should apply confirmed edit plans")
    assert_true("plan?.mode !== \"advisory\"" in js, "UI should only expose apply controls for non-advisory pending actions")
    assert_true("/api/git/diff" in js and "/api/git/restore" in js, "UI should expose Git diff and restore workflows")
    assert_true("rawDownload && downloadPercent !== null" in js, "download progress should display the current file percentage from task payloads")
    assert_true("downloadedSize" in js and "totalSize" in js, "download progress should show downloaded and total file size")
    assert_true("function clearInvalidProjectPath" in js, "UI should clear invalid project paths after open-project failure")
    assert_true('code !== "PROJECT_PATH_INVALID"' in js, "invalid-path cleanup should only run for PROJECT_PATH_INVALID")
    assert_true("elements.projectPath.value = \"\"" in js, "invalid project path should be removed from the input")


def test_static_ui_exposes_codegraph_tools_split_scripts_and_virtual_tree():
    html = (ROOT / "webui" / "static" / "index.html").read_text(encoding="utf-8")
    css = (ROOT / "webui" / "static" / "styles.css").read_text(encoding="utf-8")
    js = "\n".join(path.read_text(encoding="utf-8") for path in sorted((ROOT / "webui" / "static" / "js").glob("app-*.js")))
    codegraph_js = ROOT / "webui" / "static" / "js" / "app-codegraph.js"
    tree_js = ROOT / "webui" / "static" / "js" / "app-tree.js"
    assert_true('id="codeGraphToolbar"' in html, "chat UI should expose a CodeGraph toolbar")
    assert_true('id="codeGraphResults"' in html, "chat UI should expose CodeGraph results")
    assert_true('id="structurePanel"' in html, "chat UI should expose project structure analysis results")
    assert_true('id="codeGraphQueryInput"' not in html, "CodeGraph should reuse the main chat input instead of a second query field")
    assert_true("查詢關聯" in html, "CodeGraph primary action should clearly describe relationship lookup")
    assert_true("分析專案檔案結構" in html, "analysis action should be positioned as file-structure analysis")
    assert_true("/js/app-state.js" in html and "/js/app-main.js" in html, "index should load split plain scripts")
    assert_true(".codegraph-toolbar" in css and ".codegraph-results" in css, "CodeGraph UI should have dedicated styles")
    assert_true(".structure-panel" in css and ".structure-grid" in css, "file-structure analysis should have dedicated layout styles")
    assert_true(codegraph_js.exists(), "CodeGraph behavior should live in a split script")
    assert_true(tree_js.exists(), "virtual file tree behavior should live in a split script")
    codegraph_source = codegraph_js.read_text(encoding="utf-8")
    assert_true("queryCodeGraph" in codegraph_source, "CodeGraph split script should implement query behavior")
    assert_true("elements.chatInput.value.trim()" in codegraph_source, "CodeGraph query should read the shared chat input")
    assert_true("renderCodeGraphOperation" in codegraph_source, "CodeGraph operations should render inline progress and results")
    assert_true("renderCodeGraphStatusPanel" in codegraph_source, "CodeGraph should render a persistent status panel")
    assert_true("codegraph-chip" in codegraph_source, "CodeGraph should expose clickable query suggestion chips")
    assert_true("codeGraphLastRebuild" in js, "CodeGraph should preserve rebuild results when a follow-up query runs")
    assert_true("codeGraphQueryInput" not in js, "split scripts should not keep stale CodeGraph query input references")
    assert_true("/api/project/structure" in js, "Analyze file structure should use the deterministic structure API")
    assert_true("pinStructureRecommendedFiles" in js, "file-structure analysis should expose one-click recommended pinning")
    assert_true("virtualTree" in tree_js.read_text(encoding="utf-8"), "tree split script should implement virtualization state")
    assert_true("function maintainChatAutoScroll" in js, "streaming chat should use a guarded auto-scroll helper")
    assert_true("state.chatScroll.followLatest" in js, "streaming chat should track whether the user wants to follow the latest output")
    assert_true("function appendLiveText" in js and "maintainChatAutoScroll();" in js, "streaming token append should preserve manual scroll position")


def main():
    tests = [
        test_no_context_chat_payload,
        test_request_max_tokens_clamps_to_default,
        test_qwen_request_options_disable_thinking,
        test_default_model_uses_last_used_preference,
        test_chat_exchange_persists_last_used_model_preference,
        test_edit_plan_flow_uses_requested_model_key,
        test_gemma_context_window_matches_local_bench,
        test_context_calibration_overrides_input_budget_and_model_payload,
        test_context_calibration_runner_targets_selected_model_and_context,
        test_context_benchmark_reuses_existing_compatible_model_server,
        test_start_server_uses_low_memory_model_env_for_qwen36_on_low_vram,
        test_gemma_manifest_uses_unsloth_with_mmproj,
        test_new_model_catalog_exposes_hardware_metadata,
        test_hardware_profile_classification_and_recommendations,
        test_qwen36a3b_catalog_targets_8gb_nvidia_moe_offload,
        test_llama_launcher_accepts_auto_hardware_args,
        test_launch_webui_restarts_stale_codeworker_server,
        test_bootstrap_stops_codeworker_runtime_users_before_winpython_update,
        test_model_download_progress_payload_reports_file_percent,
        test_partial_model_file_is_rejected_by_size_guard,
        test_hardware_optimization_log_entry_contains_diagnostics,
        test_model_file_matching_does_not_fallback_on_pattern_miss,
        test_project_structure_classifies_multi_language_files,
        test_http_error_body_is_preserved,
        test_rag_manifest_search_and_stale,
        test_code_graph_indexes_symbols_and_relationships,
        test_codeworker_codegraph_plugin_is_installable_and_queryable,
        test_codegraph_api_helpers_return_status_context_and_coverage,
        test_rag_model_loading_locator_prefers_source_chunks,
        test_rag_chinese_game_speed_query_finds_code,
        test_project_rag_context_without_pins,
        test_chat_context_sends_full_line_numbered_pinned_file_for_all_models_when_budget_allows,
        test_chat_context_falls_back_to_excerpt_for_all_models_when_full_file_exceeds_budget,
        test_edit_context_sends_full_pinned_file_when_budget_allows_existing_effect_changes,
        test_edit_context_includes_all_pinned_files_when_total_budget_allows,
        test_edit_candidate_resolution_uses_rag_and_codegraph_without_pins,
        test_project_rag_rebuilds_graphless_existing_index,
        test_generated_text_file_requires_confirmation,
        test_edit_actions_apply_with_git_checkpoint_and_restore,
        test_edit_action_security_rejects_unsafe_paths_and_stale_patches,
        test_edit_action_supports_create_replace_delete_rename_and_command,
        test_edit_apply_returns_validation_command_suggestions,
        test_advisory_edit_plan_keeps_local_context_when_model_patch_is_unsafe,
        test_precise_validation_failure_logs_raw_reply_and_returns_unverified_advisory,
        test_model_precise_patch_creates_applyable_action_without_hardcoded_rule,
        test_malformed_model_patch_is_salvaged_when_search_replace_are_unique,
        test_fallback_advisory_uses_pending_target_to_locate_real_method,
        test_fallback_advisory_salvages_partial_json_without_noisy_failure_text,
        test_edit_plan_timeout_short_circuits_to_local_fallback,
        test_timeout_fallback_can_create_applyable_tetris_clear_effect_patch,
        test_edit_plan_timeout_is_scaled_for_local_coding_models,
        test_edit_plan_model_call_uses_streaming_and_scaled_timeout,
        test_streaming_edit_plan_timeout_logs_partial_reply,
        test_local_direct_edit_changes_down_key_to_soft_drop_and_ctrl_to_hard_drop,
        test_generation_prompt_infers_multiple_documents_from_previous_answer,
        test_generation_prompt_infers_excel,
        test_generation_word_prompt_uses_previous_answer,
        test_generation_with_previous_keyword_is_not_continuation,
        test_generation_common_text_aliases,
        test_generated_docx_and_text_previews_can_be_created,
        test_inline_docx_generation_uses_pasted_content_without_model,
        test_previous_answer_docx_generation_uses_history_without_model,
        test_thread_continuation_generation_loads_requested_thread_history,
        test_thread_save_updates_default_title_and_restores_project_context,
        test_missing_thread_project_path_is_cleared_instead_of_restored,
        test_invalid_open_project_path_leaves_server_idle,
        test_cleanup_empty_threads_only_removes_e2e_threads,
        test_tool_transcript_items_do_not_enter_model_history,
        test_generic_previous_answer_file_generation_defaults_to_markdown,
        test_generation_without_project_uses_app_root_and_previous_answer,
        test_generated_pdf_keeps_chinese_text_extractable,
        test_document_generation_cleans_markdown_for_pptx,
        test_document_generation_splits_long_pptx_sections,
        test_model_initiated_generation_uses_model_title_for_filename,
        test_generation_system_prompt_is_only_added_for_generation_requests,
        test_stream_chat_initializes_model_generation_flag,
        test_static_ui_exposes_file_tree_layout_and_ai_busy_indicator,
        test_static_ui_exposes_codegraph_tools_split_scripts_and_virtual_tree,
        test_gemma_multimodal_payload_and_fallback,
        test_image_metadata_fallback_blocks_guessing,
        test_video_metadata_fallback_blocks_guessing,
        test_video_timestamp_selection_handles_short_videos,
        test_media_assessment_exposes_local_limits,
        test_transcribe_media_attachment_updates_text_preview,
        test_history_continuation_uses_previous_answer_tail,
        test_chat_messages_include_recent_history,
        test_chat_messages_include_compressed_memory_summary,
        test_compact_session_memory_keeps_ui_history_and_builds_summary,
        test_length_continuation_drops_large_project_context,
        test_partial_stream_reply_can_be_saved_for_continue,
        test_stream_reasoning_only_length_retries_for_final_answer,
        test_gemma_native_image_payload_with_mmproj,
        test_prepare_attachments_does_not_use_qwen_helper,
        test_stream_attachment_fallback_for_native_model,
    ]
    try:
        for test in tests:
            test()
            print(f"PASS {test.__name__}")
    finally:
        server.cleanup_image_upload_dir()


def live_gemma_smoke():
    base_url = "http://127.0.0.1:8764"
    if "--base-url" in sys.argv:
        index = sys.argv.index("--base-url")
        if index + 1 < len(sys.argv):
            base_url = sys.argv[index + 1].rstrip("/")
    png = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
    upload_body = json.dumps(
        {
            "name": "one.png",
            "mimeType": "image/png",
            "data": "data:image/png;base64," + png,
        }
    ).encode("utf-8")
    upload_request = urllib.request.Request(
        f"{base_url}/api/uploads/file",
        data=upload_body,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    with urllib.request.urlopen(upload_request, timeout=30) as response:
        upload = json.loads(response.read().decode("utf-8"))
    attachment_id = upload["data"]["id"]
    chat_body = json.dumps(
        {
            "message": "請只回答 OK。",
            "modelKey": "gemma4",
            "attachmentIds": [attachment_id],
            "maxTokens": 32,
        },
        ensure_ascii=False,
    ).encode("utf-8")
    chat_request = urllib.request.Request(
        f"{base_url}/api/chat/stream",
        data=chat_body,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    events = []
    event_payloads = []
    with urllib.request.urlopen(chat_request, timeout=90) as response:
        for raw_line in response:
            line = raw_line.decode("utf-8", errors="replace").strip()
            if line.startswith("event:"):
                events.append(line.split(":", 1)[1].strip())
            if line.startswith("data:"):
                event_payloads.append(line.split(":", 1)[1].strip())
            if line == "event: done":
                break
    joined_payloads = "\n".join(event_payloads)
    helper_status = "vision" + "-helper"
    assert_true("qwen35" not in joined_payloads.lower(), "live Gemma4 image smoke must not use a secondary vision model")
    assert_true(helper_status not in joined_payloads.lower(), "live Gemma4 image smoke must not emit secondary vision status")
    assert_true("done" in events, "live Gemma4 image smoke should finish")
    print("PASS live_gemma_smoke")


if __name__ == "__main__":
    if "--live-gemma" in sys.argv:
        live_gemma_smoke()
    else:
        main()
